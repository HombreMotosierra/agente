import datetime
from difflib import SequenceMatcher
import json
import random
import string
import queue
import re
import sqlite3
import subprocess
import threading
import os
import shutil
import sys
import webbrowser
from pathlib import Path
import tkinter as tk
from tkinter import messagebox, ttk

import ollama
import pyttsx3
import requests
try:
    import speech_recognition as sr
except Exception:
    sr = None


# ========================
# PATHS / CONFIG
# ========================
BASE_DIR = Path(__file__).resolve().parent
APP_DIR = BASE_DIR.parent / "agente_ia_data"
APP_DIR.mkdir(exist_ok=True)
MEJORAS_DIR = APP_DIR / "mejoras_seguras"
MEJORAS_DIR.mkdir(exist_ok=True)

DB_PATH = APP_DIR / "historial.db"
CONFIG_PATH = APP_DIR / "config.json"

DEFAULT_CONFIG = {
    "model": "qwen2.5:7b",
    "voz_activa": True,
    "voz_rate": 180,
    "voz_volume": 1.0,
    "voz_voice_id": "",
    "voz_style": "Natural",
    "voz_speed_label": "Normal",
    "voz_entrada_activa": True,
    "voz_entrada_idioma": "es-ES",
    "voz_entrada_microfono": "",
    "modelo_online": "llama-3.1-8b-instant",
    "proveedor_ia": "local",
    "groq_api_key": "",
    "api_base_url": "",
    "api_bearer_token": "",
    "n8n_webhook_url": "",
    "orquestador_url": "http://127.0.0.1:8765",
    "orquestador_token": "",
    "usar_habilidades_auto": True,
    "known_paths": {},
    "window_size": "520x760",
    "compact_window_size": "340x520",
    "window_x": 1000,
    "window_y": 140,
}

DESTRUCTIVE_PATTERNS = [
    r"\bdel\b",
    r"\berase\b",
    r"\brmdir\b",
    r"\bformat\b",
    r"\bremove-item\b",
    r"\bclear-disk\b",
    r"\bdelete\b",
    r"\bdiskpart\b",
    r"rm\s+-rf",
]

NEGATIVE_PATTERNS = (
    "no puedo",
    "no sé",
    "no tengo la capacidad",
    "no puedo acceder",
    "no puedo ayudarte con eso",
)

MAX_MODELO_REINTENTOS = 3
MAX_JSON_EXTRACTION_CHARS = 20000


def normalizar_proveedor_ia(valor):
    v = str(valor or "").strip().lower()
    if v in ("api", "online", "remoto", "remota"):
        return "online"
    return "local"


def proveedor_ui_value(valor):
    return "api" if normalizar_proveedor_ia(valor) == "online" else "local"


def cargar_config():
    if not CONFIG_PATH.exists():
        guardar_config(DEFAULT_CONFIG)
        return DEFAULT_CONFIG.copy()

    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        data = json.load(f)

    cfg = DEFAULT_CONFIG.copy()
    cfg.update(data)
    cfg["proveedor_ia"] = normalizar_proveedor_ia(cfg.get("proveedor_ia"))
    return cfg


def guardar_config(config):
    with open(CONFIG_PATH, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


CONFIG = cargar_config()
ULTIMA_SOLICITUD_USUARIO = ""
escucha_activa = False
root = None
chat_window = None
style = None
estado_var = None
header = None
titulo = None
hero_subtitle_var = None
provider_badge_var = None
provider_badge = None
tts_estado = None
model_var = None
combo_model = None
provider_var = None
combo_provider = None
online_model_var = None
combo_online_model = None
mic_var = None
combo_mic = None
voz_style_var = None
combo_voz_style = None
voz_speed_var = None
combo_voz_speed = None
voz_var = None
chk_voz = None
skills_auto_var = None
chk_skills_auto = None
btn_compacto = None
btn_guardar = None
btn_habilidades = None
btn_mejoras = None
btn_automatizaciones = None
btn_integraciones = None
btn_key = None
frame_chat = None
canvas = None
scrollable_frame = None
settings_panel = None
settings_visible = True
btn_settings = None
frame_input = None
entrada = None
btn_limpiar = None
btn_micro = None
btn_refresh_mic = None
btn_enviar = None
chat_compacto = False
auto_compact_job = None


def normalizar_window_size(size):
    if not isinstance(size, str):
        return DEFAULT_CONFIG["window_size"]
    match = re.match(r"^(\d+)x(\d+)$", size.strip())
    if not match:
        return DEFAULT_CONFIG["window_size"]
    w, h = int(match.group(1)), int(match.group(2))
    if w < 320 or h < 420:
        return DEFAULT_CONFIG["window_size"]
    return f"{w}x{h}"


def _obtener_usuario_windows_preferido():
    known = CONFIG.get("known_paths", {}) if isinstance(CONFIG.get("known_paths"), dict) else {}
    user_cfg = str(known.get("windows_user", "")).strip()
    if user_cfg:
        return user_cfg
    return os.environ.get("USERNAME", "") or Path.home().name


def _guardar_known_path(clave, valor):
    if not valor:
        return
    known = CONFIG.get("known_paths")
    if not isinstance(known, dict):
        known = {}
    known[clave] = str(valor)
    CONFIG["known_paths"] = known
    guardar_config(CONFIG)


def descubrir_ruta_escritorio():
    known = CONFIG.get("known_paths", {}) if isinstance(CONFIG.get("known_paths"), dict) else {}
    ruta_cfg = known.get("desktop", "")
    if ruta_cfg:
        p = Path(ruta_cfg)
        if p.exists():
            return p

    candidatos = []
    candidatos.append(Path.home() / "Desktop")
    userprofile = os.environ.get("USERPROFILE", "").strip()
    if userprofile:
        candidatos.append(Path(userprofile) / "Desktop")

    usuario = _obtener_usuario_windows_preferido()
    letras = "CDEFGHIJKLMNOPQRSTUVWXYZ"
    for l in letras:
        root = Path(f"{l}:\\")
        if not root.exists():
            continue
        candidatos.extend(
            [
                root / "Users" / usuario / "Desktop",
                root / "Usuarios" / usuario / "Desktop",
            ]
        )

    for p in candidatos:
        if p.exists() and p.is_dir():
            _guardar_known_path("desktop", str(p))
            _guardar_known_path("windows_user", usuario)
            return p
    return Path.home() / "Desktop"


def _detectar_usuario_desde_texto(texto):
    t = (texto or "").strip()
    if not t:
        return None
    patrones = [
        r"(?i)\bmi\s+usuario\s+es\s+([a-z0-9._-]{2,64})\b",
        r"(?i)\busuario\s*:\s*([a-z0-9._-]{2,64})\b",
        r"(?i)\buser\s*:\s*([a-z0-9._-]{2,64})\b",
    ]
    for pat in patrones:
        m = re.search(pat, t)
        if m:
            return m.group(1)
    return None


def _descubrir_home_windows(usuario=None):
    usuario = (usuario or _obtener_usuario_windows_preferido() or "").strip()
    candidatos = []
    if usuario:
        for l in "CDEFGHIJKLMNOPQRSTUVWXYZ":
            root = Path(f"{l}:\\")
            if not root.exists():
                continue
            candidatos.extend(
                [
                    root / "Users" / usuario,
                    root / "Usuarios" / usuario,
                ]
            )
    userprofile = os.environ.get("USERPROFILE", "").strip()
    if userprofile:
        candidatos.append(Path(userprofile))
    candidatos.append(Path.home())

    for p in candidatos:
        if p.exists() and p.is_dir():
            _guardar_known_path("windows_user", p.name)
            _guardar_known_path("windows_home", str(p))
            return p
    return Path.home()


CONFIG["window_size"] = normalizar_window_size(CONFIG.get("window_size"))


# ========================
# DB
# ========================
def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS conversaciones (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                fecha TEXT NOT NULL,
                rol TEXT NOT NULL,
                texto TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS habilidades (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                trigger_texto TEXT NOT NULL,
                acciones_json TEXT NOT NULL,
                creado_en TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS ejecuciones (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                fecha TEXT NOT NULL,
                solicitud TEXT NOT NULL,
                resultado TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS automatizaciones (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                nombre TEXT NOT NULL,
                trigger_texto TEXT NOT NULL,
                acciones_json TEXT NOT NULL,
                habilitada INTEGER NOT NULL DEFAULT 1,
                creado_en TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS cache_consultas (
                cache_key TEXT PRIMARY KEY,
                respuesta TEXT NOT NULL,
                actualizado_en TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS lecciones (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                fecha TEXT NOT NULL,
                solicitud TEXT NOT NULL,
                acciones_json TEXT NOT NULL,
                resultado TEXT NOT NULL,
                exito INTEGER NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS mejoras_seguras (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                fecha TEXT NOT NULL,
                tipo TEXT NOT NULL,
                titulo TEXT NOT NULL,
                solicitud TEXT NOT NULL,
                detalle TEXT NOT NULL,
                propuesta_json TEXT NOT NULL,
                estado TEXT NOT NULL DEFAULT 'pendiente',
                origen TEXT NOT NULL DEFAULT 'manual',
                requiere_autorizacion INTEGER NOT NULL DEFAULT 1
            )
            """
        )
        # Migración desde esquema legado (prompt/respuesta/created_at)
        cols = [r[1] for r in conn.execute("PRAGMA table_info(conversaciones)").fetchall()]
        expected = {"fecha", "rol", "texto"}
        if cols and not expected.issubset(set(cols)):
            legacy = conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='conversaciones_legacy'"
            ).fetchone()
            if legacy is None:
                conn.execute("ALTER TABLE conversaciones RENAME TO conversaciones_legacy")
                conn.execute(
                    """
                    CREATE TABLE conversaciones (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        fecha TEXT NOT NULL,
                        rol TEXT NOT NULL,
                        texto TEXT NOT NULL
                    )
                    """
                )
                legacy_cols = [r[1] for r in conn.execute("PRAGMA table_info(conversaciones_legacy)").fetchall()]
                if {"prompt", "respuesta", "created_at"}.issubset(set(legacy_cols)):
                    rows = conn.execute(
                        "SELECT created_at, prompt, respuesta FROM conversaciones_legacy ORDER BY id ASC"
                    ).fetchall()
                    for created_at, prompt, respuesta in rows:
                        fecha = created_at or datetime.datetime.now().isoformat()
                        conn.execute(
                            "INSERT INTO conversaciones (fecha, rol, texto) VALUES (?, ?, ?)",
                            (fecha, "user", str(prompt or "")),
                        )
                        conn.execute(
                            "INSERT INTO conversaciones (fecha, rol, texto) VALUES (?, ?, ?)",
                            (fecha, "assistant", str(respuesta or "")),
                        )
        # Migración de cache_consultas desde esquemas previos
        cache_cols = [r[1] for r in conn.execute("PRAGMA table_info(cache_consultas)").fetchall()]
        if cache_cols and "actualizado_en" not in set(cache_cols):
            legacy_cache = conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' AND name='cache_consultas_legacy'"
            ).fetchone()
            if legacy_cache is None:
                conn.execute("ALTER TABLE cache_consultas RENAME TO cache_consultas_legacy")
                conn.execute(
                    """
                    CREATE TABLE cache_consultas (
                        cache_key TEXT PRIMARY KEY,
                        respuesta TEXT NOT NULL,
                        actualizado_en TEXT NOT NULL
                    )
                    """
                )
                legacy_cols = [r[1] for r in conn.execute("PRAGMA table_info(cache_consultas_legacy)").fetchall()]
                # Mapea distintos formatos históricos.
                if {"cache_key", "respuesta", "updated_at"}.issubset(set(legacy_cols)):
                    rows = conn.execute(
                        "SELECT cache_key, respuesta, updated_at FROM cache_consultas_legacy"
                    ).fetchall()
                    for cache_key, respuesta, updated_at in rows:
                        conn.execute(
                            "INSERT OR REPLACE INTO cache_consultas (cache_key, respuesta, actualizado_en) VALUES (?, ?, ?)",
                            (str(cache_key or ""), str(respuesta or ""), str(updated_at or datetime.datetime.now().isoformat())),
                        )
                elif {"clave", "respuesta", "updated_at"}.issubset(set(legacy_cols)):
                    rows = conn.execute(
                        "SELECT clave, respuesta, updated_at FROM cache_consultas_legacy"
                    ).fetchall()
                    for clave, respuesta, updated_at in rows:
                        conn.execute(
                            "INSERT OR REPLACE INTO cache_consultas (cache_key, respuesta, actualizado_en) VALUES (?, ?, ?)",
                            (str(clave or ""), str(respuesta or ""), str(updated_at or datetime.datetime.now().isoformat())),
                        )
        conn.commit()


def guardar_mensaje_db(rol, texto):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            "INSERT INTO conversaciones (fecha, rol, texto) VALUES (?, ?, ?)",
            (datetime.datetime.now().isoformat(), rol, texto),
        )
        conn.commit()


def obtener_contexto_db(limite=20):
    with sqlite3.connect(DB_PATH) as conn:
        rows = conn.execute(
            "SELECT rol, texto FROM conversaciones ORDER BY id DESC LIMIT ?",
            (limite,),
        ).fetchall()
    rows.reverse()
    return rows


def construir_resumen_contexto(memoria, max_chars=1600):
    if not memoria:
        return ""
    lineas = []
    for rol, txt in memoria:
        r = "Usuario" if str(rol).lower() == "user" else "Asistente"
        t = (txt or "").strip().replace("\n", " ")
        if len(t) > 220:
            t = t[:220] + "..."
        lineas.append(f"- {r}: {t}")
    out = "\n".join(lineas)
    if len(out) > max_chars:
        out = out[-max_chars:]
    return out


def limpiar_historial_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("DELETE FROM conversaciones")
        conn.commit()


def normalizar_texto_cache(texto):
    t = (texto or "").lower().strip()
    t = re.sub(r"[^\w\sáéíóúñü]", " ", t, flags=re.IGNORECASE)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def obtener_cache_respuesta(cache_key, max_age_seconds=60):
    with sqlite3.connect(DB_PATH) as conn:
        row = conn.execute(
            "SELECT respuesta, actualizado_en FROM cache_consultas WHERE cache_key=?",
            (cache_key,),
        ).fetchone()
    if not row:
        return None
    respuesta, actualizado_en = row
    try:
        ts = datetime.datetime.fromisoformat(actualizado_en)
    except ValueError:
        return None
    edad = (datetime.datetime.now() - ts).total_seconds()
    if edad > max_age_seconds:
        return None
    return respuesta


def guardar_cache_respuesta(cache_key, respuesta):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            INSERT INTO cache_consultas (cache_key, respuesta, actualizado_en)
            VALUES (?, ?, ?)
            ON CONFLICT(cache_key) DO UPDATE SET
                respuesta=excluded.respuesta,
                actualizado_en=excluded.actualizado_en
            """,
            (cache_key, respuesta, datetime.datetime.now().isoformat()),
        )
        conn.commit()


def guardar_habilidad(trigger_texto, acciones):
    trigger_norm = normalizar_texto_cache(trigger_texto)
    acciones_canon = json.dumps(acciones, ensure_ascii=False, sort_keys=True)

    with sqlite3.connect(DB_PATH) as conn:
        rows = conn.execute(
            "SELECT id, trigger_texto, acciones_json FROM habilidades ORDER BY id DESC LIMIT 300"
        ).fetchall()
        for hid, trigger_existente, acciones_existente in rows:
            t_exist = normalizar_texto_cache(trigger_existente)
            ratio = SequenceMatcher(None, trigger_norm, t_exist).ratio()
            try:
                parsed_exist = json.loads(acciones_existente)
                parsed_canon = json.dumps(parsed_exist, ensure_ascii=False, sort_keys=True)
            except Exception:
                parsed_canon = acciones_existente or ""
            if ratio >= 0.9 and parsed_canon == acciones_canon:
                conn.execute(
                    "UPDATE habilidades SET trigger_texto=?, creado_en=? WHERE id=?",
                    (
                        trigger_texto.strip()[:200],
                        datetime.datetime.now().isoformat(),
                        hid,
                    ),
                )
                conn.commit()
                return

        conn.execute(
            "INSERT INTO habilidades (trigger_texto, acciones_json, creado_en) VALUES (?, ?, ?)",
            (
                trigger_texto.strip()[:200],
                json.dumps(acciones, ensure_ascii=False),
                datetime.datetime.now().isoformat(),
            ),
        )
        conn.commit()


def buscar_habilidad(texto):
    consulta = normalizar_texto_cache(texto)
    if not consulta:
        return None

    with sqlite3.connect(DB_PATH) as conn:
        rows = conn.execute(
            "SELECT trigger_texto, acciones_json FROM habilidades ORDER BY id DESC LIMIT 100"
        ).fetchall()

    mejor_accion = None
    mejor_score = 0.0
    for trigger, acciones_json in rows:
        t = normalizar_texto_cache(trigger)
        if not t:
            continue
        palabras = [w for w in re.findall(r"[a-zA-Z0-9áéíóúñ]+", t) if len(w) > 2]
        if not palabras:
            continue
        score = sum(1 for w in palabras if w in consulta)
        ratio = score / max(1, len(palabras))
        # Evita disparos accidentales por coincidencias mínimas.
        if ratio >= 0.75 or t == consulta:
            try:
                return json.loads(acciones_json)
            except json.JSONDecodeError:
                return None
    return None


def respuesta_conversacional_local(texto):
    t = normalizar_texto_cache(texto)
    if not t:
        return None

    usuario = _obtener_usuario_windows_preferido() or "amigo"

    if any(k in t for k in ("sabes quien soy", "quien soy", "quién soy")):
        return f"Sí, te tengo como {usuario}."

    if any(k in t for k in ("hola", "buenas", "que tal", "qué tal")):
        return f"Hola {usuario}. ¿En qué te ayudo?"

    if "que puedes hacer" in t or "qué puedes hacer" in t:
        return (
            "Puedo ayudarte con archivos/rutas en tu equipo (listar, buscar, leer, escribir, mover, copiar), "
            "consultas de fecha y hora, abrir URLs y usar API (modo online/local). "
            "Si quieres, te doy ejemplos concretos."
        )

    if "chiste" in t and "ingenier" in t:
        return (
            "Chiste de ingenieros: \"Funciona en mi máquina\" no es una solución, "
            "es una condición de carrera."
        )

    if "chiste" in t:
        return "Aquí va uno: ¿Cuál es el colmo de un programador? Tener problemas de clase."

    return None


def asegurar_habilidades_base():
    base = [
        ("lista los archivos del escritorio", [{"accion": "listar_directorio", "args": {"path": str(descubrir_ruta_escritorio())}}]),
        ("que elementos tengo en mi escritorio", [{"accion": "listar_directorio", "args": {"path": str(descubrir_ruta_escritorio())}}]),
        ("buscar archivos *.txt en escritorio", [{"accion": "buscar_archivos", "args": {"path": str(descubrir_ruta_escritorio()), "patron": "*.txt"}}]),
        ("abrir google", [{"accion": "abrir_url", "args": {"url": "https://www.google.com"}}]),
    ]
    with sqlite3.connect(DB_PATH) as conn:
        total = conn.execute("SELECT COUNT(*) FROM habilidades").fetchone()[0]
    if total > 0:
        return
    for trigger, acciones in base:
        guardar_habilidad(trigger, acciones)


def registrar_ejecucion(solicitud, resultado):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            "INSERT INTO ejecuciones (fecha, solicitud, resultado) VALUES (?, ?, ?)",
            (
                datetime.datetime.now().isoformat(),
                solicitud[:300],
                resultado[:1500],
            ),
        )
        conn.commit()


def registrar_leccion(solicitud, acciones, resultado, exito):
    try:
        acciones_json = json.dumps(acciones or [], ensure_ascii=False)
    except TypeError:
        acciones_json = "[]"
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            INSERT INTO lecciones (fecha, solicitud, acciones_json, resultado, exito)
            VALUES (?, ?, ?, ?, ?)
            """,
            (
                datetime.datetime.now().isoformat(),
                str(solicitud or "")[:300],
                acciones_json,
                str(resultado or "")[:2000],
                1 if exito else 0,
            ),
        )
        conn.commit()


def obtener_habilidades(limite=300):
    with sqlite3.connect(DB_PATH) as conn:
        return conn.execute(
            """
            SELECT id, trigger_texto, acciones_json, creado_en
            FROM habilidades
            ORDER BY id DESC
            LIMIT ?
            """,
            (limite,),
        ).fetchall()


def actualizar_habilidad(habilidad_id, trigger_texto, acciones_json):
    trigger = (trigger_texto or "").strip()
    if not trigger:
        return False, "El trigger no puede estar vacío."

    try:
        parsed = json.loads(acciones_json)
        if not isinstance(parsed, list):
            return False, "acciones_json debe ser una lista JSON."
    except json.JSONDecodeError as ex:
        return False, f"JSON inválido: {ex}"

    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            "UPDATE habilidades SET trigger_texto=?, acciones_json=? WHERE id=?",
            (trigger[:200], json.dumps(parsed, ensure_ascii=False), habilidad_id),
        )
        conn.commit()
    return True, "Habilidad actualizada."


def eliminar_habilidad_id(habilidad_id):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("DELETE FROM habilidades WHERE id=?", (habilidad_id,))
        conn.commit()


def obtener_automatizaciones(limite=300):
    with sqlite3.connect(DB_PATH) as conn:
        return conn.execute(
            """
            SELECT id, nombre, trigger_texto, acciones_json, habilitada, creado_en
            FROM automatizaciones
            ORDER BY id DESC
            LIMIT ?
            """,
            (limite,),
        ).fetchall()


def crear_automatizacion(nombre, trigger_texto, acciones_json, habilitada=True):
    nom = (nombre or "").strip()
    trig = (trigger_texto or "").strip().lower()
    if not nom or not trig:
        return False, "Nombre y trigger son obligatorios."
    try:
        parsed = json.loads(acciones_json)
        if not isinstance(parsed, list):
            return False, "acciones_json debe ser una lista JSON."
    except json.JSONDecodeError as ex:
        return False, f"JSON inválido: {ex}"
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            INSERT INTO automatizaciones (nombre, trigger_texto, acciones_json, habilitada, creado_en)
            VALUES (?, ?, ?, ?, ?)
            """,
            (
                nom[:120],
                trig[:220],
                json.dumps(parsed, ensure_ascii=False),
                1 if habilitada else 0,
                datetime.datetime.now().isoformat(),
            ),
        )
        conn.commit()
    return True, "Automatización creada."


def actualizar_automatizacion(auto_id, nombre, trigger_texto, acciones_json, habilitada=True):
    nom = (nombre or "").strip()
    trig = (trigger_texto or "").strip().lower()
    if not nom or not trig:
        return False, "Nombre y trigger son obligatorios."
    try:
        parsed = json.loads(acciones_json)
        if not isinstance(parsed, list):
            return False, "acciones_json debe ser una lista JSON."
    except json.JSONDecodeError as ex:
        return False, f"JSON inválido: {ex}"
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            UPDATE automatizaciones
            SET nombre=?, trigger_texto=?, acciones_json=?, habilitada=?
            WHERE id=?
            """,
            (nom[:120], trig[:220], json.dumps(parsed, ensure_ascii=False), 1 if habilitada else 0, auto_id),
        )
        conn.commit()
    return True, "Automatización actualizada."


def eliminar_automatizacion(auto_id):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("DELETE FROM automatizaciones WHERE id=?", (auto_id,))
        conn.commit()


def borrar_todas_automatizaciones():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("DELETE FROM automatizaciones")
        conn.commit()


def _sanear_tipo_mejora(tipo):
    valor = normalizar_texto_cache(tipo)
    permitidos = {"habilidad", "automatizacion", "configuracion", "codigo", "prompt"}
    return valor if valor in permitidos else "codigo"


def crear_mejora_segura(tipo, titulo, solicitud, detalle, propuesta, origen="manual", requiere_autorizacion=True):
    tipo_s = _sanear_tipo_mejora(tipo)
    titulo_s = (titulo or "Mejora segura pendiente").strip()[:180]
    solicitud_s = (solicitud or "").strip()[:400]
    detalle_s = (detalle or "").strip()[:4000]
    try:
        propuesta_json = json.dumps(propuesta or {}, ensure_ascii=False, indent=2)
    except TypeError:
        propuesta_json = json.dumps({"detalle": str(propuesta or "")}, ensure_ascii=False, indent=2)

    with sqlite3.connect(DB_PATH) as conn:
        row = conn.execute(
            """
            SELECT id, propuesta_json
            FROM mejoras_seguras
            WHERE tipo=? AND solicitud=? AND estado IN ('pendiente', 'aprobada')
            ORDER BY id DESC
            LIMIT 1
            """,
            (tipo_s, solicitud_s),
        ).fetchone()
        if row:
            mejora_id = int(row[0])
            conn.execute(
                """
                UPDATE mejoras_seguras
                SET fecha=?, titulo=?, detalle=?, propuesta_json=?, origen=?, requiere_autorizacion=?
                WHERE id=?
                """,
                (
                    datetime.datetime.now().isoformat(),
                    titulo_s,
                    detalle_s,
                    propuesta_json,
                    str(origen or "manual")[:80],
                    1 if requiere_autorizacion else 0,
                    mejora_id,
                ),
            )
            conn.commit()
            return mejora_id

        cur = conn.execute(
            """
            INSERT INTO mejoras_seguras
            (fecha, tipo, titulo, solicitud, detalle, propuesta_json, estado, origen, requiere_autorizacion)
            VALUES (?, ?, ?, ?, ?, ?, 'pendiente', ?, ?)
            """,
            (
                datetime.datetime.now().isoformat(),
                tipo_s,
                titulo_s,
                solicitud_s,
                detalle_s,
                propuesta_json,
                str(origen or "manual")[:80],
                1 if requiere_autorizacion else 0,
            ),
        )
        conn.commit()
        return int(cur.lastrowid)


def obtener_mejoras_seguras(limite=300, estado=None):
    with sqlite3.connect(DB_PATH) as conn:
        if estado:
            return conn.execute(
                """
                SELECT id, tipo, titulo, solicitud, detalle, propuesta_json, estado, origen, requiere_autorizacion, fecha
                FROM mejoras_seguras
                WHERE estado=?
                ORDER BY id DESC
                LIMIT ?
                """,
                (estado, limite),
            ).fetchall()
        return conn.execute(
            """
            SELECT id, tipo, titulo, solicitud, detalle, propuesta_json, estado, origen, requiere_autorizacion, fecha
            FROM mejoras_seguras
            ORDER BY id DESC
            LIMIT ?
            """,
            (limite,),
        ).fetchall()


def obtener_mejora_segura_por_id(mejora_id):
    with sqlite3.connect(DB_PATH) as conn:
        return conn.execute(
            """
            SELECT id, tipo, titulo, solicitud, detalle, propuesta_json, estado, origen, requiere_autorizacion, fecha
            FROM mejoras_seguras
            WHERE id=?
            """,
            (mejora_id,),
        ).fetchone()


def actualizar_estado_mejora_segura(mejora_id, estado):
    estado_s = normalizar_texto_cache(estado)
    if estado_s not in {"pendiente", "aprobada", "rechazada", "aplicada"}:
        return False, "Estado de mejora inválido."
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("UPDATE mejoras_seguras SET estado=? WHERE id=?", (estado_s, mejora_id))
        conn.commit()
    return True, f"Mejora marcada como {estado_s}."


def _exportar_mejora_codigo(mejora_id, titulo, detalle, propuesta):
    fecha_txt = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r"[^a-zA-Z0-9_-]+", "_", (titulo or "mejora").strip())[:50] or "mejora"
    destino = MEJORAS_DIR / f"{fecha_txt}_{mejora_id}_{safe_name}.md"
    contenido = [
        f"# Mejora segura #{mejora_id}",
        "",
        f"Titulo: {titulo}",
        "",
        "## Detalle",
        detalle or "(sin detalle)",
        "",
        "## Propuesta JSON",
        "```json",
        json.dumps(propuesta or {}, ensure_ascii=False, indent=2),
        "```",
        "",
    ]
    destino.write_text("\n".join(contenido), encoding="utf-8")
    return destino


def aplicar_mejora_segura(mejora_id):
    fila = obtener_mejora_segura_por_id(mejora_id)
    if not fila:
        return False, "No encontré esa mejora segura."

    _id, tipo, titulo, solicitud, detalle, propuesta_json, estado, _origen, _req, _fecha = fila
    if estado == "rechazada":
        return False, "Esa mejora fue rechazada y no se puede aplicar."

    try:
        propuesta = json.loads(propuesta_json or "{}")
    except json.JSONDecodeError:
        propuesta = {}

    payload = propuesta.get("propuesta", propuesta) if isinstance(propuesta, dict) else {}

    if tipo == "habilidad":
        trigger = str(payload.get("trigger_texto") or solicitud or titulo).strip()
        acciones = payload.get("acciones", [])
        if not trigger or not isinstance(acciones, list) or not acciones:
            return False, "La mejora no trae una habilidad aplicable."
        guardar_habilidad(trigger, acciones)
        actualizar_estado_mejora_segura(mejora_id, "aplicada")
        return True, f"Habilidad aplicada: {trigger}"

    if tipo == "automatizacion":
        nombre = str(payload.get("nombre") or titulo or "Automatización segura").strip()
        trigger = str(payload.get("trigger_texto") or solicitud or titulo).strip()
        acciones = payload.get("acciones", [])
        habilitada = bool(payload.get("habilitada", True))
        if not nombre or not trigger or not isinstance(acciones, list) or not acciones:
            return False, "La mejora no trae una automatización aplicable."
        ok, msg = crear_automatizacion(nombre, trigger, json.dumps(acciones, ensure_ascii=False), habilitada=habilitada)
        if not ok:
            return False, msg
        actualizar_estado_mejora_segura(mejora_id, "aplicada")
        return True, msg

    if tipo == "configuracion":
        updates = payload.get("config_updates", {})
        if not isinstance(updates, dict) or not updates:
            return False, "La mejora no trae cambios de configuración aplicables."
        permitidas = {
            "usar_habilidades_auto",
            "proveedor_ia",
            "modelo_online",
            "model",
            "voz_activa",
            "voz_style",
            "voz_speed_label",
            "compact_window_size",
            "window_size",
        }
        aplicadas = []
        for clave, valor in updates.items():
            if clave in permitidas:
                if clave == "proveedor_ia":
                    CONFIG[clave] = normalizar_proveedor_ia(valor)
                else:
                    CONFIG[clave] = valor
                aplicadas.append(clave)
        if not aplicadas:
            return False, "No encontré cambios seguros para aplicar automáticamente."
        guardar_config(CONFIG)
        actualizar_estado_mejora_segura(mejora_id, "aplicada")
        return True, "Configuración aplicada: " + ", ".join(aplicadas)

    destino = _exportar_mejora_codigo(mejora_id, titulo, detalle, propuesta)
    actualizar_estado_mejora_segura(mejora_id, "aprobada")
    return True, f"Mejora aprobada y exportada para revisión en: {destino}"


def buscar_automatizacion_por_trigger(texto):
    t = normalizar_texto_cache(texto)
    if not t:
        return None
    autos = obtener_automatizaciones(limite=500)
    mejor = None
    mejor_score = 0.0
    for _id, nombre, trigger_texto, acciones_json, habilitada, _creado in autos:
        if not habilitada:
            continue
        trigger = normalizar_texto_cache(trigger_texto)
        if not trigger:
            continue
        seq = SequenceMatcher(None, t, trigger).ratio()
        tokens_t = {w for w in re.findall(r"[a-zA-Z0-9áéíóúñ]+", t) if len(w) > 2}
        tokens_tr = {w for w in re.findall(r"[a-zA-Z0-9áéíóúñ]+", trigger) if len(w) > 2}
        overlap = (len(tokens_t & tokens_tr) / max(1, len(tokens_tr))) if tokens_tr else 0.0
        contiene = 1.0 if (trigger in t or t in trigger) else 0.0

        # Solo considerar matches realmente cercanos para evitar desvíos.
        if not (seq >= 0.9 or overlap >= 0.9 or t == trigger):
            continue

        score = max(seq, overlap, contiene)
        try:
            acciones = json.loads(acciones_json)
            if isinstance(acciones, list) and score > mejor_score:
                mejor = {"nombre": nombre, "acciones": acciones}
                mejor_score = score
        except json.JSONDecodeError:
            continue
    return mejor


# ========================
# VOZ ROBUSTA (COLA)
# ========================
voz_queue = queue.Queue()
_voz_lock = threading.Lock()
tts_status_var = None


def set_tts_status(texto):
    if tts_status_var is not None and chat_window is not None:
        chat_window.after(0, lambda: tts_status_var.set(texto))


def listar_voces_disponibles():
    try:
        eng = pyttsx3.init()
        voces = eng.getProperty("voices") or []
        resultado = []
        for v in voces:
            nombre = getattr(v, "name", "sin_nombre")
            vid = getattr(v, "id", "")
            langs = getattr(v, "languages", [])
            langs_txt = ",".join([str(x) for x in langs]) if langs else ""
            resultado.append({"id": vid, "name": nombre, "langs": langs_txt})
        return resultado
    except Exception:
        return []


def elegir_voz_id(preferencias):
    prefs = [p.lower() for p in preferencias if p]
    voces = listar_voces_disponibles()
    if not voces:
        return ""
    for voz in voces:
        huella = f"{voz['name']} {voz['id']} {voz['langs']}".lower()
        if all(p in huella for p in prefs):
            return voz["id"]
    for voz in voces:
        huella = f"{voz['name']} {voz['id']} {voz['langs']}".lower()
        if any(p in huella for p in prefs):
            return voz["id"]
    return voces[0]["id"]


def aplicar_config_engine(engine):
    try:
        engine.setProperty("rate", int(CONFIG.get("voz_rate", 180)))
    except Exception:
        pass
    try:
        volume = float(CONFIG.get("voz_volume", 1.0))
        volume = max(0.1, min(1.0, volume))
        engine.setProperty("volume", volume)
    except Exception:
        pass
    voz_id = (CONFIG.get("voz_voice_id") or "").strip()
    if voz_id:
        try:
            engine.setProperty("voice", voz_id)
        except Exception:
            pass


def crear_engine():
    eng = pyttsx3.init()
    aplicar_config_engine(eng)
    return eng


def worker_voz():
    engine = crear_engine()
    while True:
        item = voz_queue.get()
        if isinstance(item, tuple):
            texto, evento = item
        else:
            texto, evento = item, None
        if texto is None:
            break
        if not CONFIG.get("voz_activa", True):
            if evento:
                evento.set()
            continue
        try:
            aplicar_config_engine(engine)
            engine.say(texto)
            engine.runAndWait()
        except Exception:
            try:
                engine.stop()
            except Exception:
                pass
            engine = crear_engine()
            # Reintento inmediato una vez.
            try:
                with _voz_lock:
                    set_tts_status("🔊 reintentando...")
                    aplicar_config_engine(engine)
                    engine.say(texto)
                    engine.runAndWait()
            except Exception:
                pass
        finally:
            set_tts_status("🛑 silencio")
            if evento:
                evento.set()


threading.Thread(target=worker_voz, daemon=True).start()


def hablar(texto, esperar=False):
    if not texto.strip():
        return
    voz_queue.put(texto)


def hablar_garantizado(texto):
    try:
        hablar(texto)
    except Exception:
        pass


def configurar_voz(estilo=None, velocidad=None, volumen=None, genero=None):
    cambios = []

    if velocidad is not None:
        try:
            vel = int(velocidad)
            vel = max(120, min(260, vel))
            CONFIG["voz_rate"] = vel
            cambios.append(f"velocidad={vel}")
        except (TypeError, ValueError):
            return "No pude aplicar velocidad. Usa un número entre 120 y 260."

    if volumen is not None:
        try:
            vol = float(volumen)
            if vol > 1:
                vol = vol / 100.0
            vol = max(0.1, min(1.0, vol))
            CONFIG["voz_volume"] = vol
            cambios.append(f"volumen={int(vol * 100)}%")
        except (TypeError, ValueError):
            return "No pude aplicar volumen. Usa 0.1-1.0 o 10-100."

    preferencias = []
    if genero:
        g = str(genero).lower()
        if g in ("femenina", "mujer", "female"):
            preferencias.extend(["female", "maria", "helena", "zira", "es"])
        elif g in ("masculina", "hombre", "male"):
            preferencias.extend(["male", "david", "pablo", "jorge", "es"])
    if estilo:
        e = str(estilo).lower()
        if "suave" in e or "calida" in e or "cálida" in e:
            preferencias.extend(["female", "es"])
            CONFIG["voz_style"] = "Suave"
        if "profunda" in e or "grave" in e:
            preferencias.extend(["male", "es"])
            CONFIG["voz_style"] = "Profunda"
        if "natural" in e:
            CONFIG["voz_style"] = "Natural"

    if preferencias:
        voice_id = elegir_voz_id(preferencias)
        if voice_id:
            CONFIG["voz_voice_id"] = voice_id
            cambios.append("voz=actualizada")

    guardar_config(CONFIG)
    if not cambios:
        return "No hubo cambios de voz para aplicar."
    return "Voz actualizada en tiempo real: " + ", ".join(cambios)


def aplicar_ajuste_voz_ui(style_label=None, speed_label=None):
    args = {}
    if style_label:
        style_map = {
            "Natural": "natural",
            "Suave": "suave",
            "Profunda": "grave",
            "Femenina": "suave",
            "Masculina": "grave",
        }
        estilo = style_map.get(style_label, "natural")
        args["estilo"] = estilo
        if style_label == "Femenina":
            args["genero"] = "femenina"
        elif style_label == "Masculina":
            args["genero"] = "masculina"
        CONFIG["voz_style"] = style_label

    if speed_label:
        speed_map = {"Lenta": 150, "Normal": 180, "Rápida": 220}
        vel = speed_map.get(speed_label, 180)
        args["velocidad"] = vel
        CONFIG["voz_speed_label"] = speed_label

    if not args:
        return "Sin cambios de voz."
    return configurar_voz(**args)


def listar_microfonos():
    if sr is None:
        return []
    try:
        return sr.Microphone.list_microphone_names()
    except Exception:
        return []


# ========================
# HERRAMIENTAS
# ========================
def normalizar_ruta(path):
    if path is None:
        return Path.home()
    p = str(path).strip().strip("\"'").rstrip("}").rstrip("*").strip()
    p_low = p.lower().replace("\\", "/")
    home_win = _descubrir_home_windows()
    desktop_win = descubrir_ruta_escritorio()
    if "/home/usuario" in p_low or "/home/user" in p_low:
        p = re.sub(r"(?i)/home/(usuario|user)", str(home_win).replace("\\", "/"), p)
    if "/home/" in p_low and ("/desktop" in p_low or "/escritorio" in p_low):
        p = re.sub(r"(?i)/home/[^/]+/(desktop|escritorio)", str(desktop_win).replace("\\", "/"), p)
    if "c:/users/usuario" in p_low or "c:\\users\\usuario" in p.lower():
        p = re.sub(r"(?i)c:[\\/]+users[\\/]+usuario", str(home_win), p)
    if "/path/escritorio" in p_low or "\\path\\escritorio" in p_low:
        p = str(desktop_win)
    if p_low.endswith("/escritorio") or p_low.endswith("\\escritorio"):
        p = str(desktop_win)
    if p.lower().endswith("\\temp") and ("appdata\\local\\temp" in p.lower()):
        return Path(os.environ.get("TEMP", p))
    return Path(p)


def obtener_hora():
    return str(datetime.datetime.now())


def detectar_consulta_fecha_hora(texto):
    t = (texto or "").lower()
    pide_hora = (
        "hora" in t
        or "qué hora" in t
        or "que hora" in t
        or "time" in t
    )
    pide_fecha = (
        "fecha" in t
        or "día" in t
        or "dia" in t
        or "hoy es" in t
        or "date" in t
    )
    if pide_hora and pide_fecha:
        return "fecha_hora"
    if pide_hora:
        return "hora"
    if pide_fecha:
        return "fecha"
    return None


def formatear_fecha_hora_bonita(tipo):
    ahora = datetime.datetime.now()
    dias = [
        "lunes",
        "martes",
        "miércoles",
        "jueves",
        "viernes",
        "sábado",
        "domingo",
    ]
    meses = [
        "enero",
        "febrero",
        "marzo",
        "abril",
        "mayo",
        "junio",
        "julio",
        "agosto",
        "septiembre",
        "octubre",
        "noviembre",
        "diciembre",
    ]
    dia_nombre = dias[ahora.weekday()]
    mes_nombre = meses[ahora.month - 1]
    hora_12 = ahora.strftime("%I:%M:%S")
    ampm = "a. m." if ahora.hour < 12 else "p. m."
    fecha_txt = f"{dia_nombre.capitalize()}, {ahora.day} de {mes_nombre} de {ahora.year}"
    hora_txt = f"{hora_12} {ampm}"

    if tipo == "fecha":
        return f"Hoy es {fecha_txt}."
    if tipo == "hora":
        return f"La hora actual es {hora_txt}."
    return f"Hoy es {fecha_txt} y la hora actual es {hora_txt}."


def detectar_intencion_principal(texto):
    t = normalizar_texto_cache(texto)
    tipo_tiempo = detectar_consulta_fecha_hora(texto)
    if tipo_tiempo:
        return "tiempo"
    if any(k in t for k in ("escritorio", "desktop", "carpeta", "directorio", "archivo", "archivos", "elementos")):
        return "archivos"
    if any(k in t for k in ("api", "token", "tokens", "limite", "límite", "cuota")):
        return "api"
    if any(k in t for k in ("ip", "red", "wifi", "ethernet")):
        return "red"
    if any(k in t for k in ("voz", "microfono", "micrófono", "mic", "habla", "tono")):
        return "voz"
    return "general"


def acciones_permitidas_por_intencion(intencion):
    mapa = {
        "tiempo": {"obtener_hora"},
        "archivos": {
            "listar_directorio",
            "leer_archivo",
            "escribir_archivo",
            "crear_archivo_especifico",
            "editar_archivo",
            "crear_carpeta",
            "eliminar_ruta",
            "mover_ruta",
            "copiar_ruta",
            "buscar_archivos",
            "abrir_ruta",
            "organizar_carpeta",
            "vaciar_carpeta",
        },
        "api": {"llamar_api", "disparar_webhook", "llamar_orquestador"},
        "red": {"obtener_ip_local", "obtener_ip_publica"},
        "voz": {"configurar_voz"},
        "general": set(),
    }
    return mapa.get(intencion, set())


def filtrar_acciones_por_intencion(acciones, texto):
    if not isinstance(acciones, list):
        return []
    intencion = detectar_intencion_principal(texto)
    # Para consultas conversacionales no ejecutamos herramientas.
    if intencion == "general":
        return []
    permitidas = acciones_permitidas_por_intencion(intencion)
    if not permitidas:
        return []
    filtradas = []
    for a in acciones:
        if not isinstance(a, dict):
            continue
        nombre = str(a.get("accion", "")).strip()
        if nombre in permitidas:
            filtradas.append(a)
    return filtradas


def obtener_ip_local():
    return subprocess.getoutput("ipconfig")


def obtener_ip_publica():
    try:
        return requests.get("https://api.ipify.org", timeout=8).text
    except requests.RequestException as ex:
        return f"No disponible: {ex}"


def listar_directorio(path=None):
    destino = normalizar_ruta(path) if path else descubrir_ruta_escritorio()
    try:
        if not destino.exists():
            return f"Ruta no encontrada: {destino}"
        if not destino.is_dir():
            return f"La ruta no es carpeta: {destino}"
        items = sorted(os.listdir(destino))
        if not items:
            return f"La carpeta está vacía: {destino}"
        vista = "\n".join(items[:120])
        return f"Contenido de {destino}:\n{vista}"
    except OSError as ex:
        return f"No se pudo listar {destino}: {ex}"


def leer_archivo(path):
    destino = normalizar_ruta(path)
    try:
        if not destino.exists():
            return f"Ruta no encontrada: {destino}"
        if destino.is_dir():
            return f"La ruta es carpeta, no archivo: {destino}"
        with open(destino, "r", encoding="utf-8", errors="replace") as f:
            contenido = f.read(8000)
        return f"Contenido de {destino}:\n{contenido}"
    except OSError as ex:
        return f"No se pudo leer {destino}: {ex}"


def escribir_archivo(path, contenido="", append=False):
    destino = normalizar_ruta(path)
    try:
        destino.parent.mkdir(parents=True, exist_ok=True)
        modo = "a" if append else "w"
        with open(destino, modo, encoding="utf-8") as f:
            f.write(contenido)
        return f"Archivo guardado: {destino}"
    except OSError as ex:
        return f"No se pudo escribir {destino}: {ex}"


def _asegurar_paquete_python(modulo, paquete_pip=None):
    try:
        __import__(modulo)
        return True, ""
    except ImportError:
        paquete = paquete_pip or modulo
        try:
            subprocess.check_output(
                [sys.executable, "-m", "pip", "install", paquete, "--quiet"],
                stderr=subprocess.STDOUT,
                timeout=120,
                text=True,
            )
            __import__(modulo)
            return True, f"Dependencia instalada: {paquete}"
        except Exception as ex:
            return False, f"No pude instalar {paquete}: {ex}"


def crear_archivo_especifico(path, tipo="txt", contenido=""):
    destino = normalizar_ruta(path)
    tipo_norm = (tipo or destino.suffix.lstrip(".") or "txt").lower()
    if not destino.suffix:
        destino = destino.with_suffix("." + tipo_norm)
    try:
        destino.parent.mkdir(parents=True, exist_ok=True)
    except OSError as ex:
        return f"No se pudo preparar la carpeta para {destino}: {ex}"

    if tipo_norm in ("txt", "md", "json", "csv", "log", "ini", "py", "js", "html", "css"):
        return escribir_archivo(str(destino), contenido, append=False)

    if tipo_norm == "docx":
        ok, info = _asegurar_paquete_python("docx", "python-docx")
        if not ok:
            return info
        import docx  # type: ignore
        doc = docx.Document()
        doc.add_paragraph(str(contenido or "Documento generado por el agente."))
        doc.save(str(destino))
        return f"Documento Word creado: {destino}"

    if tipo_norm == "xlsx":
        ok, info = _asegurar_paquete_python("openpyxl")
        if not ok:
            return info
        from openpyxl import Workbook  # type: ignore
        wb = Workbook()
        ws = wb.active
        texto = str(contenido or "valor")
        filas = [f for f in texto.splitlines() if f.strip()]
        if not filas:
            filas = [texto]
        for i, fila in enumerate(filas, start=1):
            celdas = [c.strip() for c in fila.split(",")] if "," in fila else [fila]
            for j, celda in enumerate(celdas, start=1):
                ws.cell(row=i, column=j, value=celda)
        wb.save(str(destino))
        return f"Documento Excel creado: {destino}"

    if tipo_norm == "pptx":
        ok, info = _asegurar_paquete_python("pptx", "python-pptx")
        if not ok:
            return info
        from pptx import Presentation  # type: ignore
        prs = Presentation()
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = "Presentación generada"
        body = str(contenido or "Contenido generado por el agente.")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
        prs.save(str(destino))
        return f"Presentación PowerPoint creada: {destino}"

    return f"Tipo de archivo no soportado aún: {tipo_norm}. Usa txt, docx, xlsx o pptx."


def editar_archivo(path, contenido="", reemplazar=False):
    destino = normalizar_ruta(path)
    if not destino.exists():
        return f"Ruta no encontrada: {destino}"
    if destino.is_dir():
        return f"La ruta es carpeta, no archivo: {destino}"

    ext = destino.suffix.lower().lstrip(".")
    if ext in ("docx", "xlsx", "pptx"):
        return (
            f"Edición directa de .{ext} aún no soportada sin instrucciones estructuradas. "
            "Pídeme crear una nueva versión del archivo o exportar a txt/csv."
        )
    return escribir_archivo(str(destino), contenido, append=not bool(reemplazar))


def crear_carpeta(path):
    destino = normalizar_ruta(path)
    try:
        destino.mkdir(parents=True, exist_ok=True)
        _guardar_known_path("last_folder", str(destino))
        return f"Carpeta creada/lista: {destino}"
    except OSError as ex:
        return f"No se pudo crear carpeta {destino}: {ex}"


def eliminar_ruta(path):
    destino = normalizar_ruta(path)
    if not destino.exists():
        return f"No existe: {destino}"
    if not confirmar_comando_destructivo(f"eliminar {destino}"):
        return "Eliminación cancelada por el usuario."
    try:
        if destino.is_dir():
            shutil.rmtree(destino)
            return f"Carpeta eliminada: {destino}"
        destino.unlink()
        return f"Archivo eliminado: {destino}"
    except OSError as ex:
        return f"No se pudo eliminar {destino}: {ex}"


def vaciar_carpeta(path):
    base = normalizar_ruta(path)
    if not base.exists() or not base.is_dir():
        return f"No existe carpeta: {base}"
    if not confirmar_comando_destructivo(f"vaciar {base}"):
        return "Vaciado cancelado por el usuario."
    eliminados = 0
    errores = 0
    for item in base.iterdir():
        try:
            if item.is_dir():
                shutil.rmtree(item)
            else:
                item.unlink()
            eliminados += 1
        except OSError:
            errores += 1
    return f"Vaciado completado en {base}. Eliminados={eliminados}, errores={errores}."


def mover_ruta(origen, destino):
    src = normalizar_ruta(origen)
    dst = normalizar_ruta(destino)
    try:
        if not src.exists():
            return f"No existe origen: {src}"
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dst))
        return f"Movido: {src} -> {dst}"
    except OSError as ex:
        return f"No se pudo mover: {ex}"


def copiar_ruta(origen, destino):
    src = normalizar_ruta(origen)
    dst = normalizar_ruta(destino)
    try:
        if not src.exists():
            return f"No existe origen: {src}"
        dst.parent.mkdir(parents=True, exist_ok=True)
        if src.is_dir():
            shutil.copytree(src, dst, dirs_exist_ok=True)
        else:
            shutil.copy2(src, dst)
        return f"Copiado: {src} -> {dst}"
    except OSError as ex:
        return f"No se pudo copiar: {ex}"


def buscar_archivos(path=None, patron="*"):
    base = normalizar_ruta(path) if path else Path.home()
    try:
        if not base.exists() or not base.is_dir():
            return f"Ruta base inválida: {base}"
        encontrados = [str(p) for p in base.rglob(patron)]
        if not encontrados:
            return f"Sin resultados para '{patron}' en {base}"
        vista = "\n".join(encontrados[:120])
        return f"Resultados en {base} para '{patron}':\n{vista}"
    except OSError as ex:
        return f"No se pudo buscar en {base}: {ex}"


def crear_archivos_aleatorios(path, cantidad=10, prefijo="archivo", extension=".txt", longitud=8):
    base = normalizar_ruta(path)
    try:
        base.mkdir(parents=True, exist_ok=True)
        cant = int(cantidad)
        cant = max(1, min(cant, 1000))
        ext = str(extension or ".txt").strip()
        if not ext.startswith("."):
            ext = "." + ext
        ext = ext[:10]
        longitud = int(longitud)
        longitud = max(4, min(longitud, 24))

        creados = []
        alfabeto = string.ascii_lowercase + string.digits
        for i in range(cant):
            suf = "".join(random.choice(alfabeto) for _ in range(longitud))
            nombre = f"{prefijo}_{i+1}_{suf}{ext}"
            destino = base / nombre
            with open(destino, "w", encoding="utf-8") as f:
                f.write(f"archivo generado automáticamente: {nombre}\n")
            creados.append(nombre)
        preview = "\n".join(creados[:40])
        return f"Creados {len(creados)} archivos en {base}:\n{preview}"
    except Exception as ex:
        return f"No se pudieron crear archivos en {base}: {ex}"


def abrir_ruta(path):
    destino = normalizar_ruta(path)
    if not destino.exists():
        return f"Ruta no encontrada: {destino}"
    try:
        os.startfile(str(destino))
        return f"Abierto: {destino}"
    except OSError as ex:
        return f"No se pudo abrir {destino}: {ex}"


def abrir_url(url):
    try:
        webbrowser.open(url, new=2)
        return f"URL abierta: {url}"
    except Exception as ex:
        return f"No se pudo abrir URL: {ex}"


def llamar_api(url, metodo="GET", json_data=None, data_text=None, headers=None, timeout=30):
    target = (url or "").strip()
    if not target:
        base = (CONFIG.get("api_base_url") or "").strip()
        if not base:
            return "Error: define url o api_base_url."
        target = base
    if not target.startswith("http://") and not target.startswith("https://"):
        return f"Error: URL inválida: {target}"

    method = (metodo or "GET").upper()
    hdrs = {"Content-Type": "application/json"}
    token = (CONFIG.get("api_bearer_token") or "").strip()
    if token:
        hdrs["Authorization"] = f"Bearer {token}"
    if isinstance(headers, dict):
        for k, v in headers.items():
            hdrs[str(k)] = str(v)

    if "api.example.com" in target.lower():
        return (
            "No tengo configurada tu API real todavía. "
            "En Integraciones define 'API Base URL' y, si aplica, token Bearer."
        )

    try:
        if data_text and json_data is None:
            resp = requests.request(method, target, headers=hdrs, data=str(data_text), timeout=int(timeout))
        else:
            resp = requests.request(method, target, headers=hdrs, json=json_data, timeout=int(timeout))
        txt = resp.text[:5000]
        return f"API {method} {target} -> {resp.status_code}\n{txt}"
    except requests.RequestException as ex:
        return (
            "No pude consultar la API en este momento. "
            "Revisa URL/token en Integraciones y tu conexión de red."
        )


def disparar_webhook(nombre_evento, payload=None):
    hook = (CONFIG.get("n8n_webhook_url") or "").strip()
    if not hook:
        return "Error: falta n8n_webhook_url en configuración."
    data = {"evento": nombre_evento, "payload": payload or {}, "origen": "agente_local"}
    return llamar_api(hook, metodo="POST", json_data=data)


def llamar_orquestador(accion, payload=None):
    base = (CONFIG.get("orquestador_url") or "").strip()
    if not base:
        return "Error: falta orquestador_url en configuración."
    target = base.rstrip("/") + "/accion"
    headers = {"Content-Type": "application/json"}
    token = (CONFIG.get("orquestador_token") or "").strip()
    if token:
        headers["X-Orquestador-Token"] = token
    body = {"accion": accion, "payload": payload or {}}
    try:
        resp = requests.post(target, headers=headers, json=body, timeout=30)
        txt = resp.text[:5000]
        return f"Orquestador {target} -> {resp.status_code}\n{txt}"
    except requests.RequestException as ex:
        return f"Error llamando orquestador: {ex}"


def organizar_carpeta(path=None, excluir_ext=".exe"):
    base = Path(path) if path else Path.home() / "Desktop"
    try:
        if not base.exists() or not base.is_dir():
            return f"Ruta inválida: {base}"

        categorias = {
            "Imagenes": {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"},
            "Documentos": {".pdf", ".doc", ".docx", ".txt", ".md", ".xlsx", ".pptx"},
            "Comprimidos": {".zip", ".rar", ".7z", ".tar", ".gz"},
            "Codigo": {".py", ".js", ".ts", ".java", ".cs", ".cpp", ".go", ".html", ".css"},
            "Instaladores": {".msi", ".apk", ".dmg"},
            "Otros": set(),
        }

        movidos = []
        for item in base.iterdir():
            if item.is_dir():
                continue

            ext = item.suffix.lower()
            if excluir_ext and ext == excluir_ext.lower():
                continue

            categoria = "Otros"
            for nombre_cat, exts in categorias.items():
                if exts and ext in exts:
                    categoria = nombre_cat
                    break

            destino_dir = base / categoria
            destino_dir.mkdir(exist_ok=True)
            destino = destino_dir / item.name
            if destino.exists():
                destino = destino_dir / f"{item.stem}_{int(datetime.datetime.now().timestamp())}{item.suffix}"

            shutil.move(str(item), str(destino))
            movidos.append(f"{item.name} -> {categoria}")

        if not movidos:
            return f"No hubo archivos para organizar en {base}."
        resumen = "\n".join(movidos[:120])
        return f"Organización completa en {base} (excluyendo {excluir_ext}).\n{resumen}"
    except OSError as ex:
        return f"No se pudo organizar la carpeta: {ex}"


def comando_requiere_confirmacion(cmd):
    cmd_limpio = (cmd or "").strip().lower()
    return any(re.search(pattern, cmd_limpio) for pattern in DESTRUCTIVE_PATTERNS)


def confirmar_comando_destructivo(cmd):
    if chat_window is None:
        return False
    resultado = {"ok": False}
    evento = threading.Event()

    def _ask():
        confirmado = messagebox.askyesno(
            "Confirmación requerida",
            (
                "Este comando parece destructivo y podría eliminar datos:\n\n"
                f"{cmd}\n\n"
                "¿Deseas ejecutarlo?"
            ),
            parent=chat_window,
        )
        resultado["ok"] = bool(confirmado)
        evento.set()

    chat_window.after(0, _ask)
    evento.wait()
    return resultado["ok"]


def ejecutar_cmd(cmd=None):
    if not cmd:
        return "Error: faltó parámetro cmd."

    if comando_requiere_confirmacion(cmd):
        if not confirmar_comando_destructivo(cmd):
            return "Comando cancelado por el usuario."

    try:
        out = subprocess.check_output(
            cmd,
            shell=True,
            stderr=subprocess.STDOUT,
            timeout=20,
            text=True,
        )
        return out.strip() or "(sin salida)"
    except subprocess.TimeoutExpired:
        return "Error: el comando excedió el tiempo máximo (20s)."
    except subprocess.CalledProcessError as ex:
        return f"Error ejecutando comando: {ex.output}"


# ========================
# IA
# ========================
def limpiar_groq_api_key(raw):
    texto = str(raw or "")
    texto = (
        texto.replace("\ufeff", "")
        .replace("\u200b", "")
        .replace("\r", "")
        .replace("\n", " ")
        .strip()
    )
    if not texto:
        return ""
    texto = texto.strip("\"'").strip()
    if texto.lower().startswith("bearer "):
        texto = texto[7:].strip()
    match = re.search(r"gsk_[A-Za-z0-9._-]+", texto)
    if match:
        return match.group(0)
    if "=" in texto and texto.count("=") == 1:
        _, right = texto.split("=", 1)
        texto = right.strip().strip("\"'")
    return texto.strip()


def obtener_groq_api_key(api_key_override=None):
    candidato = api_key_override
    if candidato is None:
        candidato = (CONFIG.get("groq_api_key") or "") or os.environ.get("GROQ_API_KEY", "")
    api_key = limpiar_groq_api_key(candidato)
    if not api_key:
        raise RuntimeError("Falta GROQ_API_KEY. Configúrala en la app o variable de entorno.")
    if not api_key.startswith("gsk_"):
        raise RuntimeError("La API key de Groq no tiene formato válido (debe iniciar con gsk_).")
    return api_key


def validar_groq_api_key(api_key_override=None, timeout=20):
    try:
        api_key = obtener_groq_api_key(api_key_override)
    except Exception as ex:
        return False, str(ex), ""

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    try:
        r = requests.get("https://api.groq.com/openai/v1/models", headers=headers, timeout=timeout)
    except requests.RequestException as ex:
        return False, f"No se pudo conectar a Groq: {ex}", api_key

    if r.status_code == 200:
        try:
            data = r.json()
            total = len(data.get("data", []))
            return True, f"Key válida. Groq respondió correctamente ({total} modelos visibles).", api_key
        except Exception:
            return True, "Key válida. Groq respondió correctamente.", api_key

    detalle = ""
    try:
        detalle = r.json().get("error", {}).get("message", "") or r.text[:240]
    except Exception:
        detalle = r.text[:240]

    if r.status_code == 401:
        return (
            False,
            "Groq rechazó la clave con 401 Invalid API Key. La app sí la leyó, pero esa key no está activa o no corresponde a Groq.",
            api_key,
        )

    return False, f"Groq {r.status_code}: {detalle}", api_key


def extraer_json_desde_texto(texto):
    bruto = str(texto or "").strip()
    if not bruto:
        return None
    candidatos = [bruto]

    bloques = re.findall(r"```(?:json)?\s*(.*?)```", bruto, flags=re.DOTALL | re.IGNORECASE)
    candidatos.extend([b.strip() for b in bloques if b and b.strip()])

    for open_ch, close_ch in (("{", "}"), ("[", "]")):
        start = bruto.find(open_ch)
        end = bruto.rfind(close_ch)
        if start != -1 and end != -1 and end > start:
            candidatos.append(bruto[start : end + 1].strip())

    for candidato in candidatos:
        if not candidato:
            continue
        try:
            return json.loads(candidato)
        except json.JSONDecodeError:
            continue
    return None


def _llamar_modelo_bruto(prompt_sistema, prompt_usuario):
    proveedor = normalizar_proveedor_ia(CONFIG.get("proveedor_ia", "local"))
    if proveedor == "online":
        return chat_online_groq(prompt_sistema, prompt_usuario)

    validar_ollama_disponible()
    res = ollama.chat(
        model=CONFIG["model"],
        messages=[
            {"role": "system", "content": prompt_sistema},
            {"role": "user", "content": prompt_usuario},
        ],
        keep_alive="30m",
    )
    return res["message"]["content"]


def consultar_modelo_con_reintentos(
    prompt_sistema,
    prompt_usuario,
    expect_json=False,
    max_intentos=MAX_MODELO_REINTENTOS,
):
    ultimo_error = "sin detalles"
    ultimo_texto = ""

    for intento in range(1, max_intentos + 1):
        prompt_actual = prompt_usuario
        if intento > 1:
            prompt_actual += (
                "\n\nREINTENTO OBLIGATORIO:\n"
                f"- Intento actual: {intento} de {max_intentos}.\n"
                f"- Problema del intento anterior: {ultimo_error}.\n"
                "- Corrige la salida. No te rindas. Si no puedes completar todo, entrega el mejor siguiente intento útil.\n"
            )

        try:
            contenido = _llamar_modelo_bruto(prompt_sistema, prompt_actual)
            ultimo_texto = str(contenido or "").strip()
            if not expect_json:
                if ultimo_texto:
                    return ultimo_texto
                ultimo_error = "respuesta vacía"
                continue

            data = extraer_json_desde_texto(ultimo_texto[:MAX_JSON_EXTRACTION_CHARS])
            if isinstance(data, dict):
                return data
            ultimo_error = "la respuesta no contenía un JSON válido"
        except Exception as ex:
            ultimo_error = str(ex)

    if expect_json:
        raise RuntimeError(f"No pude obtener JSON útil tras {max_intentos} intentos: {ultimo_error}")
    raise RuntimeError(f"No pude obtener respuesta útil tras {max_intentos} intentos: {ultimo_error}")


def construir_respuesta_respaldo(solicitud, motivo="", mejora_id=None):
    partes = [
        "Lo intenté varias veces y no logré resolverlo por completo todavía.",
        "No me voy a quedar mudo: puedo volver a intentarlo si me das un poco más de contexto o una ruta más concreta.",
    ]
    if motivo:
        partes.append(f"Detalle del bloqueo: {motivo}")
    if mejora_id is not None:
        partes.append(f"Además te dejé una mejora segura pendiente para revisión con ID #{mejora_id}.")
    return " ".join(partes)


def es_solicitud_auto_mejora(texto):
    t = normalizar_texto_cache(texto)
    patrones = (
        "mejorate",
        "mejora tu",
        "auto mejora",
        "automejora",
        "quiero que aprendas",
        "aprende una nueva capacidad",
        "hazte capaz",
        "agrega la capacidad",
        "mejora tu interfaz",
        "mejora tu codigo",
    )
    return any(p in t for p in patrones)


def construir_propuesta_mejora_local(solicitud, motivo="manual"):
    t = normalizar_texto_cache(solicitud)
    if any(k in t for k in ("interfaz", "ui", "ventana", "diseño")):
        tipo = "codigo"
        titulo = "Mejorar interfaz del agente"
        detalle = (
            "El agente detectó una solicitud de mejora visual. Esta mejora requiere cambios en la UI y debe "
            "quedar pendiente hasta autorización explícita."
        )
        propuesta = {
            "tipo": tipo,
            "propuesta": {
                "modulos_sugeridos": ["Codigo/texto.py"],
                "enfoque": [
                    "ajustar layout visual",
                    "mejorar jerarquía de controles",
                    "mantener selector manual local/api",
                ],
                "pruebas_sugeridas": ["abrir UI", "probar cambios de proveedor", "verificar paneles secundarios"],
                "motivo": motivo,
            },
        }
        return tipo, titulo, detalle, propuesta

    if any(k in t for k in ("automatiza", "automatizacion", "automatización")):
        tipo = "automatizacion"
        titulo = "Nueva automatización segura"
        detalle = "La solicitud parece poder resolverse con una automatización declarativa y segura."
        propuesta = {
            "tipo": tipo,
            "propuesta": {
                "nombre": "Automatización propuesta",
                "trigger_texto": solicitud[:200],
                "acciones": [],
                "habilitada": True,
                "motivo": motivo,
            },
        }
        return tipo, titulo, detalle, propuesta

    tipo = "codigo"
    titulo = "Nueva capacidad pendiente de mejora"
    detalle = (
        "La solicitud requiere una mejora más profunda del agente. Se registra como propuesta segura para revisión "
        "y no se tocará código automáticamente sin autorización."
    )
    propuesta = {
        "tipo": tipo,
        "propuesta": {
            "modulos_sugeridos": ["Codigo/texto.py", "Codigo/orquestador_local.py"],
            "objetivo": solicitud[:300],
            "motivo": motivo,
            "pruebas_sugeridas": ["probar proveedor local", "probar proveedor api", "probar flujo principal"],
        },
    }
    return tipo, titulo, detalle, propuesta


def generar_propuesta_mejora_segura(solicitud, motivo="manual"):
    memoria = obtener_contexto_db(limite=12)
    historial = construir_resumen_contexto(memoria, max_chars=900)
    prompt = f"""
Eres un arquitecto de mejoras seguras para un agente de escritorio.
Debes responder SOLO JSON válido con este formato:
{{
  "tipo": "habilidad|automatizacion|configuracion|codigo|prompt",
  "titulo": "resumen corto",
  "detalle": "explicación breve para el usuario",
  "requiere_autorizacion": true,
  "propuesta": {{
    "trigger_texto": "opcional",
    "acciones": [],
    "nombre": "opcional",
    "config_updates": {{}},
    "modulos_sugeridos": [],
    "pruebas_sugeridas": [],
    "motivo": "{motivo}"
  }}
}}

Reglas:
- Prefiere `habilidad`, `automatizacion` o `configuracion` si se puede mejorar sin tocar código.
- Usa `codigo` solo si realmente requiere cambios internos.
- Nunca propongas autoeditar código en este paso.
- Si no hay solución directa, devuelve una propuesta útil igualmente.

Historial:
{historial}

Solicitud de mejora:
{solicitud}
"""
    try:
        data = consultar_modelo_con_reintentos(
            "Diseñas mejoras seguras y autorizables para un agente local. Responde solo JSON válido.",
            prompt,
            expect_json=True,
        )
        tipo = _sanear_tipo_mejora(data.get("tipo"))
        tipo, titulo, detalle, propuesta = normalizar_propuesta_mejora_modelo(tipo, data, solicitud, motivo)
        if not detalle:
            detalle = "Se generó una propuesta de mejora segura pendiente de autorización."
        return tipo, titulo, detalle, propuesta
    except Exception:
        return construir_propuesta_mejora_local(solicitud, motivo=motivo)


def crear_mejora_segura_desde_solicitud(solicitud, motivo="manual", origen="manual"):
    tipo, titulo, detalle, propuesta = generar_propuesta_mejora_segura(solicitud, motivo=motivo)
    mejora_id = crear_mejora_segura(
        tipo=tipo,
        titulo=titulo,
        solicitud=solicitud,
        detalle=detalle,
        propuesta=propuesta,
        origen=origen,
        requiere_autorizacion=True,
    )
    return mejora_id, tipo, titulo, detalle, propuesta


def normalizar_propuesta_mejora_modelo(tipo, data, solicitud, motivo):
    if not isinstance(data, dict):
        return construir_propuesta_mejora_local(solicitud, motivo=motivo)

    propuesta = data.get("propuesta", {})
    if not isinstance(propuesta, dict):
        propuesta = {}
        data["propuesta"] = propuesta

    if tipo in {"habilidad", "automatizacion"}:
        acciones = propuesta.get("acciones", [])
        if not isinstance(acciones, list) or not acciones:
            return construir_propuesta_mejora_local(solicitud, motivo=motivo)
        for accion in acciones:
            if not isinstance(accion, dict) or "accion" not in accion:
                return construir_propuesta_mejora_local(solicitud, motivo=motivo)

    if tipo == "configuracion":
        updates = propuesta.get("config_updates", {})
        if not isinstance(updates, dict):
            return construir_propuesta_mejora_local(solicitud, motivo=motivo)
        permitidas = {
            "usar_habilidades_auto",
            "proveedor_ia",
            "modelo_online",
            "model",
            "voz_activa",
            "voz_style",
            "voz_speed_label",
            "compact_window_size",
            "window_size",
        }
        updates_filtradas = {k: v for k, v in updates.items() if k in permitidas}
        if not updates_filtradas:
            return construir_propuesta_mejora_local(solicitud, motivo=motivo)
        propuesta["config_updates"] = updates_filtradas

    return tipo, str(data.get("titulo") or "Mejora segura pendiente").strip(), str(data.get("detalle") or "").strip(), data


def chat_online_groq(prompt_sistema, prompt_usuario):
    api_key = obtener_groq_api_key()

    model_online = CONFIG.get("modelo_online", "llama-3.1-8b-instant")
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": model_online,
        "messages": [
            {"role": "system", "content": prompt_sistema},
            {"role": "user", "content": prompt_usuario},
        ],
        "temperature": 0.2,
    }
    try:
        r = requests.post(url, headers=headers, json=payload, timeout=45)
        if r.status_code >= 400:
            body = ""
            try:
                body = r.json().get("error", {}).get("message", "") or r.text[:240]
            except Exception:
                body = r.text[:240]
            raise RuntimeError(f"Groq {r.status_code}: {body}")
        data = r.json()
    except requests.RequestException as ex:
        raise RuntimeError(f"No se pudo conectar a Groq: {ex}") from ex
    return data["choices"][0]["message"]["content"]


def validar_ollama_disponible():
    base_url = (CONFIG.get("ollama_url") or "http://127.0.0.1:11434").rstrip("/")
    try:
        r = requests.get(f"{base_url}/api/tags", timeout=4)
        if r.status_code >= 400:
            raise RuntimeError(f"Ollama respondió {r.status_code}")
    except requests.RequestException as ex:
        raise RuntimeError(
            "Ollama local no está activo. Inícialo con: ollama serve"
        ) from ex


def decidir(prompt):
    memoria = obtener_contexto_db(limite=20)
    historial = construir_resumen_contexto(memoria, max_chars=1800)

    contexto = f"""
Eres un agente conectado a herramientas reales.

Herramientas disponibles:
- obtener_hora
- obtener_ip_local
- obtener_ip_publica
- listar_directorio(path)
- leer_archivo(path)
- escribir_archivo(path, contenido, append)
- crear_archivo_especifico(path, tipo, contenido)
- editar_archivo(path, contenido, reemplazar)
- crear_carpeta(path)
- eliminar_ruta(path)
- mover_ruta(origen, destino)
- copiar_ruta(origen, destino)
- buscar_archivos(path, patron)
- crear_archivos_aleatorios(path, cantidad, prefijo, extension)
- abrir_ruta(path)
- abrir_url(url)
- llamar_api(url, metodo, json_data, data_text, headers, timeout)
- disparar_webhook(nombre_evento, payload)
- llamar_orquestador(accion, payload)
- organizar_carpeta(path, excluir_ext)
- configurar_voz(estilo, velocidad, volumen, genero)
- ejecutar_cmd(cmd)

    Reglas estrictas:
- Responde SIEMPRE JSON válido.
- No inventes acciones.
- Si necesitas herramienta, agrega acciones.
- Si no necesitas herramienta, responde directo con acciones=[].
- Puedes usar comandos del sistema con ejecutar_cmd(cmd).
- Si el comando es destructivo (borrado/format), el sistema pedirá confirmación humana.
    - Nunca respondas "no puedo" o "no sé"; si hay duda, propone acciones de intento.
    - Mantén continuidad: usa el historial para recordar preferencias y contexto reciente.

Formato de salida:
{{
  "respuesta": "texto al usuario",
  "acciones": [
    {{"accion": "nombre_accion", "args": {{}}}}
  ]
}}

Historial reciente:
{historial}

Usuario:
{prompt}
"""

    data = consultar_modelo_con_reintentos(
        "Eres un asistente técnico preciso. Debes responder en JSON válido.",
        contexto,
        expect_json=True,
    )
    return json.dumps(data, ensure_ascii=False)


def ejecutar_acciones(lista):
    resultados = []
    for item in lista:
        accion = item.get("accion")
        args = item.get("args", {})

        if accion == "obtener_hora":
            r = obtener_hora()
        elif accion == "obtener_ip_local":
            r = obtener_ip_local()
        elif accion == "obtener_ip_publica":
            r = obtener_ip_publica()
        elif accion == "listar_directorio":
            r = listar_directorio(**args)
        elif accion == "leer_archivo":
            r = leer_archivo(**args)
        elif accion == "escribir_archivo":
            r = escribir_archivo(**args)
        elif accion == "crear_archivo_especifico":
            r = crear_archivo_especifico(**args)
        elif accion == "editar_archivo":
            r = editar_archivo(**args)
        elif accion == "crear_carpeta":
            r = crear_carpeta(**args)
        elif accion == "eliminar_ruta":
            r = eliminar_ruta(**args)
        elif accion == "vaciar_carpeta":
            r = vaciar_carpeta(**args)
        elif accion == "mover_ruta":
            r = mover_ruta(**args)
        elif accion == "copiar_ruta":
            r = copiar_ruta(**args)
        elif accion == "buscar_archivos":
            r = buscar_archivos(**args)
        elif accion == "crear_archivos_aleatorios":
            r = crear_archivos_aleatorios(**args)
        elif accion == "abrir_ruta":
            r = abrir_ruta(**args)
        elif accion == "abrir_url":
            r = abrir_url(**args)
        elif accion == "llamar_api":
            r = llamar_api(**args)
        elif accion == "disparar_webhook":
            r = disparar_webhook(**args)
        elif accion == "llamar_orquestador":
            r = llamar_orquestador(**args)
        elif accion == "organizar_carpeta":
            r = organizar_carpeta(**args)
        elif accion == "configurar_voz":
            r = configurar_voz(**args)
        elif accion == "ejecutar_cmd":
            r = ejecutar_cmd(**args)
        else:
            r = f"Acción no válida: {accion}"

        resultados.append(r)
    return resultados


# ========================
# UI - VENTANA MINIMALISTA
# ========================
def cerrar_app():
    if chat_window is None:
        return
    try:
        CONFIG["window_size"] = chat_window.geometry().split("+")[0]
        CONFIG["window_x"] = chat_window.winfo_x()
        CONFIG["window_y"] = chat_window.winfo_y()
        guardar_config(CONFIG)
    finally:
        chat_window.destroy()


# La interfaz se crea de forma explícita en configurar_interfaz().


def on_voice_ui_change(_event=None):
    msg = aplicar_ajuste_voz_ui(
        style_label=voz_style_var.get(),
        speed_label=voz_speed_var.get(),
    )
    estado_var.set("Voz actualizada")
    guardar_config(CONFIG)
    chat_window.after(1200, lambda: estado_var.set("Listo"))
    return msg


def on_provider_change(_event=None):
    CONFIG["proveedor_ia"] = normalizar_proveedor_ia(provider_var.get())
    CONFIG["modelo_online"] = online_model_var.get()
    guardar_config(CONFIG)
    actualizar_resumen_visual()
    estado_var.set(f"Proveedor fijo: {proveedor_ui_value(CONFIG['proveedor_ia'])}")
    chat_window.after(1200, lambda: estado_var.set("Listo"))


def on_toggle_voz():
    CONFIG["voz_activa"] = bool(voz_var.get())
    guardar_config(CONFIG)
    if CONFIG["voz_activa"]:
        set_tts_status("🛑 silencio")
    else:
        set_tts_status("🔇 voz desactivada")


def alternar_compacto():
    global chat_compacto
    chat_compacto = not chat_compacto
    if chat_compacto:
        chat_window.geometry(
            f"{CONFIG['compact_window_size']}+{chat_window.winfo_x()}+{chat_window.winfo_y()}"
        )
        btn_compacto.configure(text="Expandir")
    else:
        chat_window.geometry(
            f"{CONFIG['window_size']}+{chat_window.winfo_x()}+{chat_window.winfo_y()}"
        )
        btn_compacto.configure(text="Compacto")

def guardar_preferencias():
    CONFIG["model"] = model_var.get()
    CONFIG["proveedor_ia"] = normalizar_proveedor_ia(provider_var.get())
    CONFIG["modelo_online"] = online_model_var.get()
    CONFIG["voz_activa"] = bool(voz_var.get())
    CONFIG["usar_habilidades_auto"] = bool(skills_auto_var.get())
    CONFIG["window_size"] = chat_window.geometry().split("+")[0]
    CONFIG["window_x"] = chat_window.winfo_x()
    CONFIG["window_y"] = chat_window.winfo_y()
    guardar_config(CONFIG)
    estado_var.set("Guardado")
    if not CONFIG["voz_activa"]:
        set_tts_status("🔇 voz desactivada")
    else:
        set_tts_status("🛑 silencio")
    chat_window.after(1200, lambda: estado_var.set("Listo"))


def abrir_panel_habilidades():
    panel = tk.Toplevel(chat_window)
    panel.title("Habilidades aprendidas")
    panel.geometry("880x560")
    panel.configure(bg="#030712")

    top = tk.Frame(panel, bg="#0f172a")
    top.pack(fill="x")

    tk.Label(
        top,
        text="Panel de habilidades aprendidas",
        bg="#0f172a",
        fg="#e2e8f0",
        font=("Segoe UI Semibold", 11),
    ).pack(side="left", padx=10, pady=8)

    body = tk.Frame(panel, bg="#030712")
    body.pack(fill="both", expand=True, padx=10, pady=10)

    columns = ("id", "trigger", "creado")
    tree = ttk.Treeview(body, columns=columns, show="headings", height=14)
    tree.heading("id", text="ID")
    tree.heading("trigger", text="Trigger")
    tree.heading("creado", text="Creado")
    tree.column("id", width=60, anchor="center")
    tree.column("trigger", width=480, anchor="w")
    tree.column("creado", width=220, anchor="w")
    tree.pack(fill="x")

    tk.Label(
        body,
        text="Trigger",
        bg="#030712",
        fg="#cbd5e1",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(10, 2))
    trigger_entry = tk.Entry(
        body, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat"
    )
    trigger_entry.pack(fill="x", ipady=5)

    tk.Label(
        body,
        text="Acciones JSON",
        bg="#030712",
        fg="#cbd5e1",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(10, 2))
    acciones_txt = tk.Text(
        body,
        bg="#0b1220",
        fg="#e2e8f0",
        insertbackground="#e2e8f0",
        relief="flat",
        height=12,
    )
    acciones_txt.pack(fill="both", expand=True)

    selected_id = {"value": None}

    def cargar_tabla():
        for item in tree.get_children():
            tree.delete(item)
        for hid, trigger, _acciones, creado in obtener_habilidades():
            tree.insert("", "end", values=(hid, trigger, creado))

    def on_select(_evt=None):
        sel = tree.selection()
        if not sel:
            return
        vals = tree.item(sel[0], "values")
        if not vals:
            return
        hid = int(vals[0])
        selected_id["value"] = hid
        filas = [r for r in obtener_habilidades() if int(r[0]) == hid]
        if not filas:
            return
        _, trigger, acciones_json, _ = filas[0]
        trigger_entry.delete(0, tk.END)
        trigger_entry.insert(0, trigger)
        acciones_txt.delete("1.0", tk.END)
        acciones_txt.insert("1.0", acciones_json)

    tree.bind("<<TreeviewSelect>>", on_select)

    acciones_botones = tk.Frame(body, bg="#030712")
    acciones_botones.pack(fill="x", pady=(8, 0))

    def guardar_edicion():
        if not selected_id["value"]:
            messagebox.showinfo("Editar habilidad", "Selecciona una habilidad primero.", parent=panel)
            return
        ok, msg = actualizar_habilidad(
            selected_id["value"],
            trigger_entry.get(),
            acciones_txt.get("1.0", tk.END).strip(),
        )
        if not ok:
            messagebox.showerror("Editar habilidad", msg, parent=panel)
            return
        messagebox.showinfo("Editar habilidad", msg, parent=panel)
        cargar_tabla()

    def borrar_habilidad():
        if not selected_id["value"]:
            messagebox.showinfo("Eliminar habilidad", "Selecciona una habilidad primero.", parent=panel)
            return
        if not messagebox.askyesno(
            "Eliminar habilidad",
            "¿Seguro que deseas eliminar esta habilidad?",
            parent=panel,
        ):
            return
        eliminar_habilidad_id(selected_id["value"])
        selected_id["value"] = None
        trigger_entry.delete(0, tk.END)
        acciones_txt.delete("1.0", tk.END)
        cargar_tabla()

    tk.Button(
        acciones_botones,
        text="Recargar",
        command=cargar_tabla,
        bg="#111827",
        fg="#cbd5e1",
        relief="flat",
        padx=10,
    ).pack(side="left")
    tk.Button(
        acciones_botones,
        text="Guardar cambios",
        command=guardar_edicion,
        bg="#1d4ed8",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left", padx=6)
    tk.Button(
        acciones_botones,
        text="Eliminar",
        command=borrar_habilidad,
        bg="#7f1d1d",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left")
    tk.Button(
        acciones_botones,
        text="Eliminar todas",
        command=lambda: borrar_todas_habilidades(),
        bg="#991b1b",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left", padx=6)

    def borrar_todas_habilidades():
        if not messagebox.askyesno(
            "Eliminar todas",
            "¿Seguro que deseas eliminar TODAS las habilidades aprendidas?",
            parent=panel,
        ):
            return
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute("DELETE FROM habilidades")
            conn.commit()
        selected_id["value"] = None
        trigger_entry.delete(0, tk.END)
        acciones_txt.delete("1.0", tk.END)
        cargar_tabla()

    cargar_tabla()


def abrir_panel_mejoras():
    panel = tk.Toplevel(chat_window)
    panel.title("Mejoras seguras")
    panel.geometry("1040x640")
    panel.configure(bg="#030712")

    top = tk.Frame(panel, bg="#0f172a")
    top.pack(fill="x")
    tk.Label(
        top,
        text="Mejoras seguras pendientes y autorizadas",
        bg="#0f172a",
        fg="#e2e8f0",
        font=("Segoe UI Semibold", 11),
    ).pack(side="left", padx=10, pady=8)

    body = tk.Frame(panel, bg="#030712")
    body.pack(fill="both", expand=True, padx=10, pady=10)

    left = tk.Frame(body, bg="#030712")
    left.pack(side="left", fill="both", expand=True, padx=(0, 8))

    columns = ("id", "tipo", "titulo", "estado", "origen", "fecha")
    tree = ttk.Treeview(left, columns=columns, show="headings", height=18)
    for col, label, width in (
        ("id", "ID", 60),
        ("tipo", "Tipo", 110),
        ("titulo", "Título", 340),
        ("estado", "Estado", 100),
        ("origen", "Origen", 100),
        ("fecha", "Fecha", 180),
    ):
        tree.heading(col, text=label)
        tree.column(col, width=width, anchor="w")
    tree.pack(fill="both", expand=True)

    right = tk.Frame(body, bg="#030712")
    right.pack(side="right", fill="both", expand=True)

    resumen_var = tk.StringVar(value="Selecciona una mejora para ver su detalle.")
    tk.Label(
        right,
        textvariable=resumen_var,
        bg="#030712",
        fg="#e2e8f0",
        justify="left",
        wraplength=430,
        font=("Segoe UI Semibold", 10),
    ).pack(anchor="w")

    detalle_var = tk.StringVar(value="")
    tk.Label(
        right,
        textvariable=detalle_var,
        bg="#030712",
        fg="#94a3b8",
        justify="left",
        wraplength=430,
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(8, 10))

    tk.Label(
        right,
        text="Propuesta",
        bg="#030712",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(anchor="w")
    propuesta_txt = tk.Text(
        right,
        bg="#0b1220",
        fg="#e2e8f0",
        insertbackground="#e2e8f0",
        relief="flat",
        height=20,
    )
    propuesta_txt.pack(fill="both", expand=True)

    selected = {"id": None}

    def cargar_tabla():
        for item in tree.get_children():
            tree.delete(item)
        for mid, tipo, titulo, _sol, _det, _prop, estado, origen, _req, fecha in obtener_mejoras_seguras():
            tree.insert("", "end", values=(mid, tipo, titulo, estado, origen, fecha))

    def cargar_detalle(mejora_id):
        fila = obtener_mejora_segura_por_id(mejora_id)
        if not fila:
            resumen_var.set("No encontré esa mejora.")
            detalle_var.set("")
            propuesta_txt.delete("1.0", tk.END)
            return
        _id, tipo, titulo, solicitud, detalle, propuesta_json, estado, origen, requiere_autorizacion, fecha = fila
        resumen_var.set(f"#{_id} · {titulo} · {tipo} · {estado}")
        detalle_var.set(
            f"Solicitud: {solicitud}\nOrigen: {origen}\nFecha: {fecha}\nRequiere autorización: {'sí' if requiere_autorizacion else 'no'}\n\n{detalle}"
        )
        propuesta_txt.delete("1.0", tk.END)
        propuesta_txt.insert("1.0", propuesta_json)

    def on_select(_evt=None):
        sel = tree.selection()
        if not sel:
            return
        vals = tree.item(sel[0], "values")
        if not vals:
            return
        selected["id"] = int(vals[0])
        cargar_detalle(selected["id"])

    tree.bind("<<TreeviewSelect>>", on_select)

    acciones = tk.Frame(right, bg="#030712")
    acciones.pack(fill="x", pady=(10, 0))

    def aprobar_aplicar():
        if not selected["id"]:
            messagebox.showinfo("Mejoras", "Selecciona una mejora primero.", parent=panel)
            return
        ok, msg = aplicar_mejora_segura(selected["id"])
        if not ok:
            messagebox.showerror("Mejoras", msg, parent=panel)
            return
        sincronizar_ui_desde_config()
        messagebox.showinfo("Mejoras", msg, parent=panel)
        cargar_tabla()
        cargar_detalle(selected["id"])

    def rechazar_mejora():
        if not selected["id"]:
            messagebox.showinfo("Mejoras", "Selecciona una mejora primero.", parent=panel)
            return
        ok, msg = actualizar_estado_mejora_segura(selected["id"], "rechazada")
        if not ok:
            messagebox.showerror("Mejoras", msg, parent=panel)
            return
        messagebox.showinfo("Mejoras", msg, parent=panel)
        cargar_tabla()
        cargar_detalle(selected["id"])

    tk.Button(
        acciones,
        text="Recargar",
        command=cargar_tabla,
        bg="#111827",
        fg="#cbd5e1",
        relief="flat",
        padx=10,
    ).pack(side="left")
    tk.Button(
        acciones,
        text="Aprobar / Aplicar",
        command=aprobar_aplicar,
        bg="#1d4ed8",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left", padx=6)
    tk.Button(
        acciones,
        text="Rechazar",
        command=rechazar_mejora,
        bg="#7f1d1d",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left")

    cargar_tabla()


def pedir_api_key_groq():
    panel = tk.Toplevel(chat_window)
    panel.title("Configurar Groq API Key")
    panel.geometry("620x250")
    panel.configure(bg="#030712")

    tk.Label(
        panel,
        text="Pega tu API key de Groq:",
        bg="#030712",
        fg="#e2e8f0",
        font=("Segoe UI", 10),
    ).pack(anchor="w", padx=12, pady=(12, 6))

    tk.Label(
        panel,
        text="Acepta claves con formato `gsk_...`, incluso si pegas `Bearer ...` o `GROQ_API_KEY=...`.",
        bg="#030712",
        fg="#94a3b8",
        font=("Segoe UI", 9),
    ).pack(anchor="w", padx=12, pady=(0, 8))

    key_entry = tk.Entry(
        panel,
        show="*",
        bg="#111827",
        fg="#e2e8f0",
        insertbackground="#e2e8f0",
        relief="flat",
        font=("Segoe UI", 10),
    )
    key_entry.pack(fill="x", padx=12, ipady=5)
    key_entry.insert(0, CONFIG.get("groq_api_key", ""))

    status_var_panel = tk.StringVar(value="Pega la key y usa 'Probar key' antes de guardarla.")
    status_lbl = tk.Label(
        panel,
        textvariable=status_var_panel,
        bg="#030712",
        fg="#94a3b8",
        justify="left",
        wraplength=590,
        font=("Segoe UI", 9),
    )
    status_lbl.pack(anchor="w", padx=12, pady=(10, 4))

    def set_panel_status(texto, ok=None):
        status_var_panel.set(texto)
        if ok is True:
            status_lbl.configure(fg="#86efac")
        elif ok is False:
            status_lbl.configure(fg="#fca5a5")
        else:
            status_lbl.configure(fg="#94a3b8")

    def alternar_visible():
        if key_entry.cget("show") == "*":
            key_entry.configure(show="")
            btn_toggle.configure(text="Ocultar")
        else:
            key_entry.configure(show="*")
            btn_toggle.configure(text="Mostrar")

    def pegar_portapapeles():
        try:
            raw_clip = panel.clipboard_get()
        except tk.TclError:
            set_panel_status("No encontré texto en el portapapeles.", ok=False)
            return
        key_clean = limpiar_groq_api_key(raw_clip)
        key_entry.delete(0, tk.END)
        key_entry.insert(0, key_clean)
        if key_clean.startswith("gsk_"):
            set_panel_status("Key detectada en el portapapeles. Ya puedes probarla.", ok=True)
        else:
            set_panel_status("No encontré una key de Groq válida en el texto pegado.", ok=False)

    def probar_key():
        raw = key_entry.get() or ""
        key_clean = limpiar_groq_api_key(raw)
        key_entry.delete(0, tk.END)
        key_entry.insert(0, key_clean)
        set_panel_status("Probando conectividad con Groq...", ok=None)
        panel.update_idletasks()
        ok, detalle, key_validada = validar_groq_api_key(key_clean)
        if ok:
            set_panel_status(detalle, ok=True)
            return True, key_validada
        set_panel_status(detalle, ok=False)
        return False, key_validada

    def guardar_key():
        ok, key_validada = probar_key()
        if not ok:
            messagebox.showerror(
                "Groq API Key",
                "La key no se guardó porque Groq la rechazó o no tiene formato correcto.\n\n"
                + status_var_panel.get(),
                parent=panel,
            )
            estado_var.set("Error en key Groq")
            chat_window.after(1800, lambda: estado_var.set("Listo"))
            return

        CONFIG["groq_api_key"] = key_validada
        guardar_config(CONFIG)
        actualizar_resumen_visual()
        estado_var.set("Groq conectado")
        chat_window.after(1400, lambda: estado_var.set("Listo"))
        messagebox.showinfo("Groq API Key", "Key validada y guardada correctamente.", parent=panel)
        panel.destroy()

    acciones = tk.Frame(panel, bg="#030712")
    acciones.pack(fill="x", padx=12, pady=10)

    btn_pegar = tk.Button(
        acciones,
        text="Pegar",
        command=pegar_portapapeles,
        bg="#111827",
        fg="#cbd5e1",
        relief="flat",
        padx=12,
    )
    btn_pegar.pack(side="left")

    btn_toggle = tk.Button(
        acciones,
        text="Mostrar",
        command=alternar_visible,
        bg="#111827",
        fg="#cbd5e1",
        relief="flat",
        padx=12,
    )
    btn_toggle.pack(side="left", padx=(8, 0))

    btn_probar = tk.Button(
        acciones,
        text="Probar key",
        command=probar_key,
        bg="#0f766e",
        fg="#f8fafc",
        relief="flat",
        padx=12,
    )
    btn_probar.pack(side="right")

    btn_guardar = tk.Button(
        acciones,
        text="Guardar key",
        command=guardar_key,
        bg="#1d4ed8",
        fg="#f8fafc",
        relief="flat",
        padx=12,
    )
    btn_guardar.pack(side="right", padx=(0, 8))


def abrir_panel_integraciones():
    panel = tk.Toplevel(chat_window)
    panel.title("Integraciones API / Webhook")
    panel.geometry("760x320")
    panel.configure(bg="#030712")

    body = tk.Frame(panel, bg="#030712")
    body.pack(fill="both", expand=True, padx=12, pady=12)

    tk.Label(body, text="API base URL", bg="#030712", fg="#e2e8f0", font=("Segoe UI", 10)).pack(anchor="w")
    api_url = tk.Entry(body, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat")
    api_url.pack(fill="x", ipady=5, pady=(2, 8))
    api_url.insert(0, CONFIG.get("api_base_url", ""))

    tk.Label(body, text="API Bearer token", bg="#030712", fg="#e2e8f0", font=("Segoe UI", 10)).pack(anchor="w")
    api_token = tk.Entry(
        body, show="*", bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat"
    )
    api_token.pack(fill="x", ipady=5, pady=(2, 8))
    api_token.insert(0, CONFIG.get("api_bearer_token", ""))

    tk.Label(body, text="n8n Webhook URL", bg="#030712", fg="#e2e8f0", font=("Segoe UI", 10)).pack(anchor="w")
    hook_url = tk.Entry(body, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat")
    hook_url.pack(fill="x", ipady=5, pady=(2, 12))
    hook_url.insert(0, CONFIG.get("n8n_webhook_url", ""))

    tk.Label(body, text="Orquestador URL", bg="#030712", fg="#e2e8f0", font=("Segoe UI", 10)).pack(anchor="w")
    orch_url = tk.Entry(body, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat")
    orch_url.pack(fill="x", ipady=5, pady=(2, 8))
    orch_url.insert(0, CONFIG.get("orquestador_url", "http://127.0.0.1:8765"))

    tk.Label(body, text="Orquestador token", bg="#030712", fg="#e2e8f0", font=("Segoe UI", 10)).pack(anchor="w")
    orch_token = tk.Entry(
        body, show="*", bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat"
    )
    orch_token.pack(fill="x", ipady=5, pady=(2, 12))
    orch_token.insert(0, CONFIG.get("orquestador_token", ""))

    def guardar_integraciones():
        CONFIG["api_base_url"] = api_url.get().strip()
        CONFIG["api_bearer_token"] = api_token.get().strip()
        CONFIG["n8n_webhook_url"] = hook_url.get().strip()
        CONFIG["orquestador_url"] = orch_url.get().strip()
        CONFIG["orquestador_token"] = orch_token.get().strip()
        guardar_config(CONFIG)
        estado_var.set("Integraciones guardadas")
        chat_window.after(1200, lambda: estado_var.set("Listo"))
        panel.destroy()

    def probar_api():
        res = llamar_api("", metodo="GET")
        messagebox.showinfo("Prueba API", res[:2000], parent=panel)

    def probar_webhook():
        res = disparar_webhook("prueba_manual", {"ok": True})
        messagebox.showinfo("Prueba Webhook", res[:2000], parent=panel)

    def probar_orquestador():
        res = llamar_orquestador("health", {})
        messagebox.showinfo("Prueba Orquestador", res[:2000], parent=panel)

    foot = tk.Frame(body, bg="#030712")
    foot.pack(fill="x")
    tk.Button(
        foot, text="Probar API", command=probar_api, bg="#111827", fg="#cbd5e1", relief="flat", padx=10
    ).pack(side="left")
    tk.Button(
        foot, text="Probar Webhook", command=probar_webhook, bg="#111827", fg="#cbd5e1", relief="flat", padx=10
    ).pack(side="left", padx=6)
    tk.Button(
        foot, text="Probar Orquestador", command=probar_orquestador, bg="#111827", fg="#cbd5e1", relief="flat", padx=10
    ).pack(side="left")
    tk.Button(
        foot, text="Guardar", command=guardar_integraciones, bg="#1d4ed8", fg="#f8fafc", relief="flat", padx=12
    ).pack(side="right")


def abrir_panel_automatizaciones():
    panel = tk.Toplevel(chat_window)
    panel.title("Automatizaciones")
    panel.geometry("980x620")
    panel.configure(bg="#030712")

    top = tk.Frame(panel, bg="#0f172a")
    top.pack(fill="x")
    tk.Label(
        top,
        text="Automatizaciones (disparadores por texto)",
        bg="#0f172a",
        fg="#e2e8f0",
        font=("Segoe UI Semibold", 11),
    ).pack(side="left", padx=10, pady=8)

    body = tk.Frame(panel, bg="#030712")
    body.pack(fill="both", expand=True, padx=10, pady=10)

    columns = ("id", "nombre", "trigger", "estado", "creado")
    tree = ttk.Treeview(body, columns=columns, show="headings", height=14)
    tree.heading("id", text="ID")
    tree.heading("nombre", text="Nombre")
    tree.heading("trigger", text="Trigger")
    tree.heading("estado", text="Estado")
    tree.heading("creado", text="Creado")
    tree.column("id", width=50, anchor="center")
    tree.column("nombre", width=170, anchor="w")
    tree.column("trigger", width=320, anchor="w")
    tree.column("estado", width=80, anchor="center")
    tree.column("creado", width=220, anchor="w")
    tree.pack(fill="x")

    form = tk.Frame(body, bg="#030712")
    form.pack(fill="both", expand=True, pady=(10, 0))

    tk.Label(form, text="Nombre", bg="#030712", fg="#cbd5e1", font=("Segoe UI", 9)).pack(anchor="w")
    nombre_entry = tk.Entry(form, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat")
    nombre_entry.pack(fill="x", ipady=5)

    tk.Label(form, text="Trigger (texto)", bg="#030712", fg="#cbd5e1", font=("Segoe UI", 9)).pack(
        anchor="w", pady=(8, 2)
    )
    trigger_entry = tk.Entry(form, bg="#111827", fg="#e2e8f0", insertbackground="#e2e8f0", relief="flat")
    trigger_entry.pack(fill="x", ipady=5)

    habilitada_var = tk.BooleanVar(value=True)
    tk.Checkbutton(
        form,
        text="Habilitada",
        variable=habilitada_var,
        bg="#030712",
        fg="#cbd5e1",
        selectcolor="#1f2937",
        activebackground="#030712",
        activeforeground="#cbd5e1",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(8, 2))

    tk.Label(form, text="Acciones JSON", bg="#030712", fg="#cbd5e1", font=("Segoe UI", 9)).pack(anchor="w")
    acciones_txt = tk.Text(
        form,
        bg="#0b1220",
        fg="#e2e8f0",
        insertbackground="#e2e8f0",
        relief="flat",
        height=12,
    )
    acciones_txt.pack(fill="both", expand=True)

    selected_id = {"value": None}

    def limpiar_form():
        selected_id["value"] = None
        nombre_entry.delete(0, tk.END)
        trigger_entry.delete(0, tk.END)
        acciones_txt.delete("1.0", tk.END)
        habilitada_var.set(True)

    def cargar_tabla():
        for item in tree.get_children():
            tree.delete(item)
        for aid, nombre, trigger, _acciones, habilitada, creado in obtener_automatizaciones():
            estado = "ON" if int(habilitada) == 1 else "OFF"
            tree.insert("", "end", values=(aid, nombre, trigger, estado, creado))

    def on_select(_evt=None):
        sel = tree.selection()
        if not sel:
            return
        vals = tree.item(sel[0], "values")
        if not vals:
            return
        aid = int(vals[0])
        filas = [r for r in obtener_automatizaciones() if int(r[0]) == aid]
        if not filas:
            return
        _id, nombre, trigger, acciones_json, habilitada, _creado = filas[0]
        selected_id["value"] = _id
        nombre_entry.delete(0, tk.END)
        nombre_entry.insert(0, nombre)
        trigger_entry.delete(0, tk.END)
        trigger_entry.insert(0, trigger)
        acciones_txt.delete("1.0", tk.END)
        acciones_txt.insert("1.0", acciones_json)
        habilitada_var.set(bool(habilitada))

    tree.bind("<<TreeviewSelect>>", on_select)

    botones = tk.Frame(form, bg="#030712")
    botones.pack(fill="x", pady=(8, 0))

    def crear_auto():
        ok, msg = crear_automatizacion(
            nombre_entry.get(),
            trigger_entry.get(),
            acciones_txt.get("1.0", tk.END).strip(),
            habilitada=habilitada_var.get(),
        )
        if not ok:
            messagebox.showerror("Automatizaciones", msg, parent=panel)
            return
        cargar_tabla()
        limpiar_form()
        messagebox.showinfo("Automatizaciones", msg, parent=panel)

    def guardar_auto():
        if not selected_id["value"]:
            messagebox.showinfo("Automatizaciones", "Selecciona una automatización primero.", parent=panel)
            return
        ok, msg = actualizar_automatizacion(
            selected_id["value"],
            nombre_entry.get(),
            trigger_entry.get(),
            acciones_txt.get("1.0", tk.END).strip(),
            habilitada=habilitada_var.get(),
        )
        if not ok:
            messagebox.showerror("Automatizaciones", msg, parent=panel)
            return
        cargar_tabla()
        messagebox.showinfo("Automatizaciones", msg, parent=panel)

    def eliminar_auto():
        if not selected_id["value"]:
            messagebox.showinfo("Automatizaciones", "Selecciona una automatización primero.", parent=panel)
            return
        if not messagebox.askyesno(
            "Automatizaciones",
            "¿Seguro que deseas eliminar esta automatización?",
            parent=panel,
        ):
            return
        eliminar_automatizacion(selected_id["value"])
        cargar_tabla()
        limpiar_form()

    def eliminar_todas():
        if not messagebox.askyesno(
            "Automatizaciones",
            "¿Eliminar TODAS las automatizaciones?",
            parent=panel,
        ):
            return
        borrar_todas_automatizaciones()
        cargar_tabla()
        limpiar_form()

    tk.Button(
        botones,
        text="Recargar",
        command=cargar_tabla,
        bg="#111827",
        fg="#cbd5e1",
        relief="flat",
        padx=10,
    ).pack(side="left")
    tk.Button(
        botones,
        text="Nueva",
        command=crear_auto,
        bg="#0f766e",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left", padx=6)
    tk.Button(
        botones,
        text="Guardar cambios",
        command=guardar_auto,
        bg="#1d4ed8",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left")
    tk.Button(
        botones,
        text="Eliminar",
        command=eliminar_auto,
        bg="#7f1d1d",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left", padx=6)
    tk.Button(
        botones,
        text="Eliminar todas",
        command=eliminar_todas,
        bg="#991b1b",
        fg="#f8fafc",
        relief="flat",
        padx=10,
    ).pack(side="left")

    cargar_tabla()

def mensaje(texto, tipo):
    if scrollable_frame is None or chat_window is None or canvas is None:
        return
    frame = tk.Frame(scrollable_frame, bg="#08111d")
    frame.pack(fill="x", padx=18, pady=5)

    if tipo == "user":
        role_txt = "Tú"
        bubble_bg = "#2563eb"
        meta_fg = "#bfdbfe"
        body_fg = "#eff6ff"
        anchor = "e"
    elif tipo == "system":
        role_txt = "Sistema"
        bubble_bg = "#1e293b"
        meta_fg = "#94a3b8"
        body_fg = "#e2e8f0"
        anchor = "w"
    else:
        role_txt = "Agente"
        bubble_bg = "#0f172a"
        meta_fg = "#7dd3fc"
        body_fg = "#e5eefb"
        anchor = "w"

    bubble_wrap = max(300, min(chat_window.winfo_width() - 170, 620))
    bubble = tk.Frame(frame, bg=bubble_bg, padx=14, pady=10)
    bubble.pack(anchor=anchor)

    meta = tk.Label(
        bubble,
        text=role_txt,
        bg=bubble_bg,
        fg=meta_fg,
        font=("Segoe UI Semibold", 8),
    )
    meta.pack(anchor="w")

    body = tk.Label(
        bubble,
        text=texto,
        bg=bubble_bg,
        fg=body_fg,
        wraplength=bubble_wrap,
        justify="left",
        font=("Segoe UI", 10),
        padx=0,
        pady=4,
    )
    body.pack(anchor="w")

    chat_window.update_idletasks()
    canvas.yview_moveto(1)


def obtener_texto_entrada():
    if entrada is None:
        return ""
    if isinstance(entrada, tk.Text):
        return entrada.get("1.0", "end-1c").strip()
    return entrada.get().strip()


def limpiar_entrada_ui():
    if entrada is None:
        return
    if isinstance(entrada, tk.Text):
        entrada.delete("1.0", tk.END)
    else:
        entrada.delete(0, tk.END)


def poner_texto_entrada(texto):
    if entrada is None:
        return
    limpiar_entrada_ui()
    if isinstance(entrada, tk.Text):
        entrada.insert("1.0", texto)
    else:
        entrada.insert(0, texto)
    entrada.focus_set()


def actualizar_resumen_visual():
    if provider_badge_var is None:
        return

    provider = normalizar_proveedor_ia(provider_var.get() if provider_var is not None else CONFIG.get("proveedor_ia"))
    provider_ui = proveedor_ui_value(provider)

    if provider == "online":
        provider_badge_var.set("API")
        if provider_badge is not None:
            provider_badge.configure(bg="#3f2a11", fg="#fde68a")
        if combo_model is not None:
            combo_model.configure(state="disabled")
        if combo_online_model is not None:
            combo_online_model.configure(state="readonly")
        detalle = f"Modo manual activado. Ahora mismo usa API con {online_model_var.get()}."
    else:
        provider_badge_var.set("LOCAL")
        if provider_badge is not None:
            provider_badge.configure(bg="#163323", fg="#86efac")
        if combo_model is not None:
            combo_model.configure(state="readonly")
        if combo_online_model is not None:
            combo_online_model.configure(state="disabled")
        detalle = f"Modo manual activado. Ahora mismo usa IA local con {model_var.get()}."

    if hero_subtitle_var is not None:
        hero_subtitle_var.set(detalle + " No cambiará de proveedor hasta que tú lo hagas.")
    if titulo is not None:
        titulo.configure(text=f"Agente IA híbrido · {provider_ui.upper()}")


def sincronizar_ui_desde_config():
    if provider_var is not None:
        provider_var.set(proveedor_ui_value(CONFIG.get("proveedor_ia", "local")))
    if model_var is not None:
        model_var.set(CONFIG.get("model", "qwen2.5:7b"))
    if online_model_var is not None:
        online_model_var.set(CONFIG.get("modelo_online", "llama-3.1-8b-instant"))
    if voz_var is not None:
        voz_var.set(bool(CONFIG.get("voz_activa", True)))
    if voz_style_var is not None:
        voz_style_var.set(CONFIG.get("voz_style", "Natural"))
    if voz_speed_var is not None:
        voz_speed_var.set(CONFIG.get("voz_speed_label", "Normal"))
    if skills_auto_var is not None:
        skills_auto_var.set(bool(CONFIG.get("usar_habilidades_auto", False)))
    if mic_var is not None:
        mic_var.set(CONFIG.get("voz_entrada_microfono", mic_var.get()))
    actualizar_resumen_visual()


def alternar_panel_ajustes():
    global settings_visible
    if settings_panel is None or btn_settings is None:
        return
    settings_visible = not settings_visible
    if settings_visible:
        settings_panel.pack(fill="x", pady=(0, 14), after=header)
        btn_settings.configure(text="Ocultar ajustes")
    else:
        settings_panel.pack_forget()
        btn_settings.configure(text="Mostrar ajustes")


def on_model_change(_event=None):
    CONFIG["model"] = model_var.get()
    guardar_config(CONFIG)
    actualizar_resumen_visual()
    estado_var.set("Modelo local actualizado")
    chat_window.after(1200, lambda: estado_var.set("Listo"))


def on_online_model_change(_event=None):
    CONFIG["modelo_online"] = online_model_var.get()
    guardar_config(CONFIG)
    actualizar_resumen_visual()
    estado_var.set("Modelo API actualizado")
    chat_window.after(1200, lambda: estado_var.set("Listo"))


def on_toggle_skills_auto():
    CONFIG["usar_habilidades_auto"] = bool(skills_auto_var.get())
    guardar_config(CONFIG)
    estado_var.set("Aprendizaje automático actualizado")
    chat_window.after(1200, lambda: estado_var.set("Listo"))


def on_textbox_return(event):
    return enviar(event)


def on_textbox_shift_return(_event):
    if entrada is None:
        return "break"
    entrada.insert(tk.INSERT, "\n")
    return "break"

def guardar_geometria_actual():
    if chat_window is None:
        return
    CONFIG["window_size"] = normalizar_window_size(chat_window.geometry().split("+")[0])
    CONFIG["window_x"] = chat_window.winfo_x()
    CONFIG["window_y"] = chat_window.winfo_y()
    guardar_config(CONFIG)


def programar_auto_compacto():
    global auto_compact_job
    if chat_window is None:
        return
    if auto_compact_job:
        chat_window.after_cancel(auto_compact_job)
    auto_compact_job = chat_window.after(18000, lambda: compactar_si_inactivo())


def compactar_si_inactivo():
    global chat_compacto, auto_compact_job
    if chat_window is None:
        return
    auto_compact_job = None
    if chat_window.focus_displayof() is None and not chat_compacto:
        chat_compacto = True
        chat_window.geometry(
            f"{CONFIG['compact_window_size']}+{chat_window.winfo_x()}+{chat_window.winfo_y()}"
        )
        btn_compacto.configure(text="Expandir")
        guardar_geometria_actual()


def expandir_si_compacto():
    global chat_compacto
    if chat_window is None:
        return
    if chat_compacto:
        chat_compacto = False
        chat_window.geometry(
            f"{CONFIG['window_size']}+{chat_window.winfo_x()}+{chat_window.winfo_y()}"
        )
        btn_compacto.configure(text="Compacto")
    programar_auto_compacto()


def limpiar_historial():
    limpiar_historial_db()
    if scrollable_frame is None:
        return
    for w in scrollable_frame.winfo_children():
        w.destroy()
    mensaje("Historial limpiado.", "system")

def resolver_ruta_alias(texto):
    t = texto.lower()
    usuario_texto = _detectar_usuario_desde_texto(texto)
    if usuario_texto:
        _guardar_known_path("windows_user", usuario_texto)
    home = _descubrir_home_windows(usuario_texto)
    escritorio_descubierto = descubrir_ruta_escritorio()
    # Tolera errores comunes de escritura: "escriitorio", "escritiro", etc.
    if "desktop" in t or "escritorio" in t or "escri" in t:
        return escritorio_descubierto
    if "documentos" in t or "documents" in t:
        return home / "Documents"
    if "descargas" in t or "downloads" in t:
        return home / "Downloads"
    if "imagenes" in t or "pictures" in t:
        return home / "Pictures"
    if "musica" in t or "music" in t:
        return home / "Music"
    if "videos" in t:
        return home / "Videos"
    if "temp" in t or "temporal" in t:
        return Path(os.environ.get("TEMP", str(home / "AppData" / "Local" / "Temp")))
    if "esa carpeta" in t or "esa ruta" in t or "ahi" in t or "ahí" in t:
        known = CONFIG.get("known_paths", {}) if isinstance(CONFIG.get("known_paths"), dict) else {}
        last_folder = str(known.get("last_folder", "")).strip()
        if last_folder:
            p = Path(last_folder)
            if p.exists() and p.is_dir():
                return p
    return None


def _obtener_ultima_carpeta_contexto():
    known = CONFIG.get("known_paths", {}) if isinstance(CONFIG.get("known_paths"), dict) else {}
    last_folder = str(known.get("last_folder", "")).strip()
    if not last_folder:
        return None
    p = Path(last_folder)
    if p.exists() and p.is_dir():
        return p
    return None


def _normalizar_nombre_carpeta_referencia(valor):
    s = (valor or "").strip()
    if not s:
        return ""
    m_nombre = re.search(r"(?i)\bcon\s+(?:el\s+)?nombre\s+([a-zA-Z0-9_\-\s]{1,80})$", s)
    if m_nombre:
        s = m_nombre.group(1).strip()
    s = re.sub(r"(?i)^\s*(la|el|una|un)\s+", "", s).strip()
    s = re.sub(r"(?i)\bcarpeta\b", "", s).strip()
    s = re.sub(r"(?i)\bllamad[ao]s?\b", "", s).strip()
    s = re.sub(r"(?i)\bde nombre\b", "", s).strip()
    s = re.sub(r"(?i)\ben\s+mi\s+escritorio\b", "", s).strip()
    s = re.sub(r"(?i)\ben\s+el\s+escritorio\b", "", s).strip()
    s = re.sub(r"(?i)\ben\s+escritorio\b", "", s).strip()
    s = re.sub(r"(?i)\bmi\s+escritorio\b", "", s).strip()
    s = re.sub(r"\s+", " ", s).strip(" .,:;")
    return s


def acciones_locales_desde_texto(texto):
    t = texto.lower().strip()
    ruta_alias = resolver_ruta_alias(t)
    usuario_texto = _detectar_usuario_desde_texto(texto)
    if usuario_texto:
        _guardar_known_path("windows_user", usuario_texto)

    if any(k in t for k in ("listar", "lista", "mostrar")) and any(
        k in t for k in ("archivos", "elementos", "carpeta", "directorio")
    ):
        destino_lista = ruta_alias or Path.home()
        if any(k in t for k in ("esa carpeta", "esa ruta", "ahí", "ahi")):
            destino_ctx = resolver_ruta_alias(t)
            if destino_ctx is not None:
                destino_lista = destino_ctx
        return [{"accion": "listar_directorio", "args": {"path": str(destino_lista)}}]

    if usuario_texto and any(k in t for k in ("escritorio", "desktop", "buscar", "encuentra")):
        return [{"accion": "listar_directorio", "args": {"path": str(descubrir_ruta_escritorio())}}]

    if ("abre " in t or "abrir " in t) and ("http://" in t or "https://" in t):
        m = re.search(r"(https?://\S+)", texto)
        if m:
            return [{"accion": "abrir_url", "args": {"url": m.group(1)}}]

    if any(k in t for k in ("llama api", "consulta api", "api get", "api post")):
        m = re.search(r"(https?://\S+)", texto)
        if m:
            metodo = "GET"
            if "post" in t:
                metodo = "POST"
            elif "put" in t:
                metodo = "PUT"
            elif "delete" in t:
                metodo = "DELETE"
            return [{"accion": "llamar_api", "args": {"url": m.group(1), "metodo": metodo}}]

    if any(k in t for k in ("dispara webhook", "lanzar webhook", "ejecuta webhook")):
        return [{"accion": "disparar_webhook", "args": {"nombre_evento": texto.strip()[:120], "payload": {}}}]

    if any(k in t for k in ("orquesta", "orquestador", "workflow")) and any(
        k in t for k in ("lanzar", "ejecuta", "dispara", "corre")
    ):
        return [
            {
                "accion": "llamar_orquestador",
                "args": {"accion": "run_workflow", "payload": {"prompt": texto.strip()}},
            }
        ]

    if ("buscar" in t or "encuentra" in t) and ("*.") in t:
        m = re.search(r"(\*\.[a-zA-Z0-9]+)", t)
        patron = m.group(1) if m else "*"
        return [{"accion": "buscar_archivos", "args": {"path": str(ruta_alias or Path.home()), "patron": patron}}]

    if (
        "carpeta" in t
        and any(k in t for k in ("crear", "crea", "nueva", "haz", "genera"))
        and not any(k in t for k in ("archivo", "archivos", "documento", "documentos"))
    ):
        m = re.search(r"['\"]([^'\"]+)['\"]", texto)
        if m:
            return [{"accion": "crear_carpeta", "args": {"path": m.group(1)}}]
        # Soporte de lenguaje natural sin comillas:
        # "crea una carpeta death en mi escritorio"
        # "en mi escritorio crea la carpeta death"
        nombre = None
        m1 = re.search(
            r"(?i)(?:crear|crea|nueva)\s+(?:una\s+)?carpeta\s+([a-zA-Z0-9_\-\s]{1,80}?)(?:\s+en\b|$)",
            texto.strip(),
        )
        if m1:
            nombre = m1.group(1).strip(" .,:;")
        if not nombre:
            m2 = re.search(
                r"(?i)carpeta\s+([a-zA-Z0-9_\-\s]{1,80}?)(?:\s+en\b|$)",
                texto.strip(),
            )
            if m2:
                nombre = m2.group(1).strip(" .,:;")
        nombre = _normalizar_nombre_carpeta_referencia(nombre)
        if nombre:
            base = ruta_alias or descubrir_ruta_escritorio()
            return [{"accion": "crear_carpeta", "args": {"path": str(base / nombre)}}]

    if (
        ("crear" in t or "crea" in t or "genera" in t)
        and ("archivo" in t or "documento" in t)
        and ("aleatorio" not in t and "aleatorios" not in t)
        and ("archivos" not in t)
        and (re.search(r"\b\d{1,4}\b", t) is None)
    ):
        m_nombre = re.search(r"['\"]([^'\"]+\.[a-zA-Z0-9]{2,5})['\"]", texto)
        m_contenido = re.search(r"(?:contenido|texto)\s*:\s*(.+)$", texto, flags=re.IGNORECASE)
        contenido = (m_contenido.group(1).strip() if m_contenido else "Generado por agente local.")
        tipo = "txt"
        if "word" in t or ".docx" in t:
            tipo = "docx"
        elif "excel" in t or ".xlsx" in t:
            tipo = "xlsx"
        elif "powerpoint" in t or "ppt" in t or ".pptx" in t:
            tipo = "pptx"
        elif ".txt" in t:
            tipo = "txt"
        if m_nombre:
            destino = normalizar_ruta(m_nombre.group(1))
        else:
            destino = (ruta_alias or descubrir_ruta_escritorio()) / f"nuevo_documento.{tipo}"
        return [
            {
                "accion": "crear_archivo_especifico",
                "args": {"path": str(destino), "tipo": tipo, "contenido": contenido},
            }
        ]

    if any(k in t for k in ("edita", "editar", "agrega al archivo", "escribe en archivo")):
        m_arch = re.search(r"['\"]([^'\"]+\.[a-zA-Z0-9]{2,5})['\"]", texto)
        m_contenido = re.search(r"(?:contenido|texto)\s*:\s*(.+)$", texto, flags=re.IGNORECASE)
        if m_arch and m_contenido:
            return [
                {
                    "accion": "editar_archivo",
                    "args": {
                        "path": m_arch.group(1),
                        "contenido": m_contenido.group(1).strip(),
                        "reemplazar": ("reemplaza" in t or "sobrescribe" in t),
                    },
                }
            ]

    if ("crear" in t or "crea" in t or "genera" in t) and ("archivo" in t or "archivos" in t or "elemento" in t or "elementos" in t):
        m_cant = re.search(r"\b(\d{1,4})\b", t)
        if m_cant is None and ("aleatorio" not in t and "aleatorios" not in t):
            # Si no pide cantidad ni aleatoriedad, no forzamos lote.
            return []
        cantidad = int(m_cant.group(1)) if m_cant else 10
        m_carpeta = re.search(
            r"(?i)(?:en|dentro de)\s+([a-zA-Z0-9_\\/:.\-\s]+?)(?:\s+(?:crea|crear|genera)\b|$)",
            texto.strip(),
        )
        destino = None
        if m_carpeta:
            ruta_txt = m_carpeta.group(1).strip()
            # "en esa carpeta" debe apuntar al último folder creado/abierto.
            if any(k in ruta_txt.lower() for k in ("esa carpeta", "esa ruta", "ahí", "ahi")):
                destino = resolver_ruta_alias(ruta_txt)
            elif "carpeta " in ruta_txt.lower():
                m_nom = re.search(r"(?i)carpeta\s+([a-zA-Z0-9_\-\s]{1,80})", ruta_txt)
                if m_nom:
                    nombre = _normalizar_nombre_carpeta_referencia(m_nom.group(1))
                    base = resolver_ruta_alias(texto) or descubrir_ruta_escritorio()
                    candidato = base / nombre
                    destino = candidato if candidato.exists() else candidato
            elif re.match(r"^[a-zA-Z0-9_\-\s]{1,80}$", ruta_txt) and ("\\" not in ruta_txt and "/" not in ruta_txt and ":" not in ruta_txt):
                nombre_ref = _normalizar_nombre_carpeta_referencia(ruta_txt)
                base = resolver_ruta_alias(texto) or descubrir_ruta_escritorio()
                destino = base / (nombre_ref or ruta_txt.strip())
            else:
                destino = normalizar_ruta(ruta_txt)
        if destino is None:
            ctx = _obtener_ultima_carpeta_contexto()
            if ctx is not None:
                destino = ctx
        if destino is None:
            destino = ruta_alias or descubrir_ruta_escritorio()
        return [
            {
                "accion": "crear_archivos_aleatorios",
                "args": {"path": str(destino), "cantidad": cantidad, "prefijo": "item", "extension": ".txt"},
            }
        ]

    if any(k in t for k in ("elimina todo", "borra todo", "vaciar", "limpia todo")) and (
        "temp" in t or "temporal" in t
    ):
        ruta_temp = Path(os.environ.get("TEMP", str(Path.home() / "AppData" / "Local" / "Temp")))
        return [{"accion": "vaciar_carpeta", "args": {"path": str(ruta_temp)}}]

    if "organiza" in t and "carpeta" in t and (
        "deja afuera el ejecutable" in t
        or "sin ejecutable" in t
        or "sin exe" in t
        or ".exe" in t
    ):
        ruta = resolver_ruta_alias(t) or Path.home() / "Desktop"
        return [
            {
                "accion": "organizar_carpeta",
                "args": {"path": str(ruta), "excluir_ext": ".exe"},
            }
        ]

    if any(k in t for k in ("voz", "tono", "habla")) and any(
        k in t for k in ("mas", "más", "pon", "cambia", "quiero")
    ):
        args = {}
        if any(k in t for k in ("femenina", "mujer")):
            args["genero"] = "femenina"
        elif any(k in t for k in ("masculina", "hombre")):
            args["genero"] = "masculina"

        if any(k in t for k in ("suave", "calida", "cálida")):
            args["estilo"] = "suave"
        elif any(k in t for k in ("grave", "profunda")):
            args["estilo"] = "grave"

        if "rápida" in t or "rapida" in t:
            args["velocidad"] = 220
        elif "lenta" in t:
            args["velocidad"] = 150

        m_vol = re.search(r"(volumen|volume)\s*(\d{1,3})", t)
        if m_vol:
            args["volumen"] = int(m_vol.group(2))

        if args:
            return [{"accion": "configurar_voz", "args": args}]

    if any(k in t for k in ("microfono", "micrófono", "mic", "escucha")) and any(
        k in t for k in ("encendido", "activar", "prender", "on")
    ):
        CONFIG["voz_entrada_activa"] = True
        guardar_config(CONFIG)
        return []

    if any(k in t for k in ("microfono", "micrófono", "mic", "escucha")) and any(
        k in t for k in ("apagar", "off", "desactivar")
    ):
        CONFIG["voz_entrada_activa"] = False
        guardar_config(CONFIG)
        return []

    return []


def es_intencion_operativa(texto):
    t = (texto or "").lower()
    claves = (
        "archivo",
        "archivos",
        "carpeta",
        "directorio",
        "ruta",
        "escritorio",
        "desktop",
        "documento",
        "office",
        "word",
        "excel",
        "powerpoint",
        "descarga",
        "listar",
        "mostrar",
        "buscar",
        "abrir",
        "mover",
        "copiar",
        "eliminar",
        "borrar",
        "vaciar",
        "api",
        "webhook",
        "orquestador",
        "workflow",
        "compila",
        "compilar",
        "instala",
        "instalar",
        "apk",
        "proyecto",
        "programa",
        "comando",
        "cmd",
    )
    return any(k in t for k in claves)


def generar_acciones_con_modelo(texto):
    memoria = obtener_contexto_db(limite=20)
    historial = construir_resumen_contexto(memoria, max_chars=1400)
    prompt_plan = f"""
Convierte la solicitud del usuario a acciones ejecutables.
Nunca respondas que no puedes. Siempre intenta una estrategia.
Devuelve JSON válido con:
{{
  "respuesta": "resumen corto",
  "acciones": [
    {{"accion": "nombre", "args": {{}}}}
  ]
}}

Acciones permitidas:
- obtener_hora
- obtener_ip_local
- obtener_ip_publica
- listar_directorio(path)
- leer_archivo(path)
- escribir_archivo(path, contenido, append)
- crear_archivo_especifico(path, tipo, contenido)
- editar_archivo(path, contenido, reemplazar)
- crear_carpeta(path)
- eliminar_ruta(path)
- mover_ruta(origen, destino)
- copiar_ruta(origen, destino)
- buscar_archivos(path, patron)
- crear_archivos_aleatorios(path, cantidad, prefijo, extension)
- abrir_ruta(path)
- abrir_url(url)
- llamar_api(url, metodo, json_data, data_text, headers, timeout)
- disparar_webhook(nombre_evento, payload)
- llamar_orquestador(accion, payload)
- organizar_carpeta(path, excluir_ext)
- vaciar_carpeta(path)
- configurar_voz(estilo, velocidad, volumen, genero)
- ejecutar_cmd(cmd)

Historial:
{historial}

Solicitud:
{texto}
"""
    try:
        data = consultar_modelo_con_reintentos(
            "Eres un planner de acciones. Devuelves JSON válido estricto.",
            prompt_plan,
            expect_json=True,
        )
        acciones = data.get("acciones", [])
        if isinstance(acciones, list):
            return data.get("respuesta", ""), acciones
    except Exception:
        pass
    return "", []


def _resultado_es_fallido(linea):
    s = (linea or "").lower()
    patrones = (
        "error",
        "no pude",
        "no tengo",
        "no encontré",
        "no encontre",
        "ruta no encontrada",
        "url inválida",
        "url invalida",
        "api.example.com",
        "invalid api key",
        "failed",
        "max retries exceeded",
    )
    return any(p in s for p in patrones)


def procesar_prompt_sync(texto):
    global ULTIMA_SOLICITUD_USUARIO

    solicitud_original = (texto or "").strip()
    if not solicitud_original:
        return "prompt vacio"

    guardar_mensaje_db("user", solicitud_original)

    texto_l = solicitud_original.lower()
    tipo_tiempo = detectar_consulta_fecha_hora(solicitud_original)
    intencion = detectar_intencion_principal(solicitud_original)

    if not any(k in texto_l for k in ("aprende", "aprendelo", "apréndelo")):
        ULTIMA_SOLICITUD_USUARIO = solicitud_original

    if tipo_tiempo:
        cache_key = f"fecha_hora::{tipo_tiempo}::{normalizar_texto_cache(solicitud_original)}"
        ttl_seg = 10 if tipo_tiempo in ("hora", "fecha_hora") else 24 * 3600
        respuesta_cache = obtener_cache_respuesta(cache_key, max_age_seconds=ttl_seg)
        respuesta_final = respuesta_cache or formatear_fecha_hora_bonita(tipo_tiempo)
        if not respuesta_cache:
            guardar_cache_respuesta(cache_key, respuesta_final)
        guardar_mensaje_db("assistant", respuesta_final)
        registrar_ejecucion(solicitud_original, respuesta_final)
        return respuesta_final

    if intencion == "general":
        resp_local = respuesta_conversacional_local(solicitud_original)
        if resp_local:
            guardar_mensaje_db("assistant", resp_local)
            registrar_ejecucion(solicitud_original, resp_local)
            return resp_local

    if es_solicitud_auto_mejora(solicitud_original):
        mejora_id, tipo, titulo_mejora, detalle_mejora, _ = crear_mejora_segura_desde_solicitud(
            solicitud_original,
            motivo="solicitud explícita del usuario",
            origen="manual",
        )
        respuesta_mejora = (
            f"Preparé una mejora segura #{mejora_id} para esto.\n\n"
            f"Título: {titulo_mejora}\n"
            f"Tipo: {tipo}\n"
            f"Detalle: {detalle_mejora}\n\n"
            "No toqué el código automáticamente. Puedes revisarla y autorizarla desde el panel `Mejoras`."
        )
        guardar_mensaje_db("assistant", respuesta_mejora)
        registrar_ejecucion(solicitud_original, respuesta_mejora)
        return respuesta_mejora

    if any(k in texto_l for k in ("aprende a hacerlo", "aprende hacerlo", "aprendelo", "apréndelo")):
        if not ULTIMA_SOLICITUD_USUARIO:
            respuesta_final = "No tengo una solicitud previa para aprender. Dame primero una tarea concreta."
        else:
            acciones_aprendidas = acciones_locales_desde_texto(ULTIMA_SOLICITUD_USUARIO)
            if not acciones_aprendidas:
                respuesta_final = (
                    "Intenté aprender, pero no encontré un patrón automático para esa tarea todavía. "
                    "Dímela con más detalle (ruta, archivo o patrón)."
                )
            else:
                guardar_habilidad(ULTIMA_SOLICITUD_USUARIO, acciones_aprendidas)
                resultados_aprendidos = ejecutar_acciones(acciones_aprendidas)
                respuesta_final = "Aprendido y ejecutado en tiempo real.\n\n" + "\n".join(resultados_aprendidos)
        guardar_mensaje_db("assistant", respuesta_final)
        registrar_ejecucion(solicitud_original, respuesta_final)
        return respuesta_final

    if CONFIG.get("usar_habilidades_auto", True):
        acciones_habilidad = buscar_habilidad(solicitud_original)
        if acciones_habilidad:
            acciones_habilidad = filtrar_acciones_por_intencion(acciones_habilidad, solicitud_original)
            if acciones_habilidad:
                resultados_habilidad = ejecutar_acciones(acciones_habilidad)
                hubo_error_habilidad = any(_resultado_es_fallido(r) for r in resultados_habilidad)
                registrar_leccion(
                    solicitud_original,
                    acciones_habilidad,
                    "\n".join(resultados_habilidad),
                    not hubo_error_habilidad,
                )
                respuesta_final = "\n".join(resultados_habilidad) or "Habilidad ejecutada."
                guardar_mensaje_db("assistant", respuesta_final)
                registrar_ejecucion(solicitud_original, respuesta_final)
                return respuesta_final

    acciones_auto = buscar_automatizacion_por_trigger(solicitud_original)
    if acciones_auto:
        acciones_auto = filtrar_acciones_por_intencion(acciones_auto, solicitud_original)
        if acciones_auto:
            resultados_auto = ejecutar_acciones(acciones_auto)
            hubo_error_auto = any(_resultado_es_fallido(r) for r in resultados_auto)
            registrar_leccion(solicitud_original, acciones_auto, "\n".join(resultados_auto), not hubo_error_auto)
            respuesta_final = "\n".join(resultados_auto) or "Automatización ejecutada."
            guardar_mensaje_db("assistant", respuesta_final)
            registrar_ejecucion(solicitud_original, respuesta_final)
            return respuesta_final

    acciones_locales = acciones_locales_desde_texto(solicitud_original) if es_intencion_operativa(solicitud_original) else []
    if acciones_locales:
        acciones_locales = filtrar_acciones_por_intencion(acciones_locales, solicitud_original)
    if acciones_locales:
        resultados_locales = ejecutar_acciones(acciones_locales)
        exito_local = not any(_resultado_es_fallido(r) for r in resultados_locales)
        registrar_leccion(solicitud_original, acciones_locales, "\n".join(resultados_locales), exito_local)
        if exito_local:
            guardar_habilidad(solicitud_original, acciones_locales)
        respuesta_final = "\n".join(resultados_locales)
        guardar_mensaje_db("assistant", respuesta_final)
        registrar_ejecucion(solicitud_original, respuesta_final)
        return respuesta_final

    if es_intencion_operativa(solicitud_original):
        propuesta_auto = []
        t_auto = solicitud_original.lower()
        base_auto = resolver_ruta_alias(solicitud_original) or descubrir_ruta_escritorio()
        if ("crear" in t_auto or "genera" in t_auto) and ("docx" in t_auto or "word" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "documento_generado.docx"), "tipo": "docx", "contenido": "Documento generado automáticamente"}}]
        elif ("crear" in t_auto or "genera" in t_auto) and ("xlsx" in t_auto or "excel" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "tabla_generada.xlsx"), "tipo": "xlsx", "contenido": "columna1,columna2\nvalor1,valor2"}}]
        elif ("crear" in t_auto or "genera" in t_auto) and ("pptx" in t_auto or "powerpoint" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "presentacion_generada.pptx"), "tipo": "pptx", "contenido": "Presentación generada automáticamente"}}]

        if propuesta_auto:
            resultados_auto = ejecutar_acciones(propuesta_auto)
            exito_auto = not any(_resultado_es_fallido(r) for r in resultados_auto)
            registrar_leccion(solicitud_original, propuesta_auto, "\n".join(resultados_auto), exito_auto)
            if exito_auto:
                guardar_habilidad(solicitud_original, propuesta_auto)
            respuesta_auto = "\n".join(resultados_auto)
            guardar_mensaje_db("assistant", respuesta_auto)
            registrar_ejecucion(solicitud_original, respuesta_auto)
            return respuesta_auto

    try:
        respuesta_cruda = decidir(solicitud_original)
    except Exception as ex:
        respuesta_cruda = json.dumps(
            {
                "respuesta": f"Error consultando proveedor IA: {ex}",
                "acciones": [],
            },
            ensure_ascii=False,
        )

    data = extraer_json_desde_texto(respuesta_cruda)
    if not isinstance(data, dict):
        texto_plano = str(respuesta_cruda or "").strip()
        data = {"respuesta": texto_plano or "No devolviste JSON útil.", "acciones": []}

    respuesta = data.get("respuesta", "").strip()
    acciones = filtrar_acciones_por_intencion(data.get("acciones", []), solicitud_original)

    if not acciones and any(
        s in respuesta.lower()
        for s in (
            "no puedo acceder",
            "no tengo la capacidad",
            "no puedo acceder directamente",
        )
    ):
        ruta_alias = resolver_ruta_alias(texto_l)
        if ruta_alias:
            acciones = [{"accion": "listar_directorio", "args": {"path": str(ruta_alias)}}]
            if not respuesta:
                respuesta = "Lo consulto localmente ahora mismo."

    requiere_plan = es_intencion_operativa(solicitud_original)
    if requiere_plan and ((not acciones) or any(p in respuesta.lower() for p in NEGATIVE_PATTERNS)):
        try:
            r_plan, a_plan = generar_acciones_con_modelo(solicitud_original)
            if a_plan:
                acciones = a_plan
                if r_plan:
                    respuesta = r_plan
            else:
                acciones = []
                if not respuesta:
                    respuesta = (
                        "No encontré un plan ejecutable para esa petición aún. "
                        "Si quieres, te muestro opciones o me das una instrucción más concreta."
                    )
        except Exception as ex:
            acciones = []
            if not respuesta:
                respuesta = f"Error en planificación: {ex}"

    if not tipo_tiempo and acciones:
        acciones = [a for a in acciones if a.get("accion") != "obtener_hora"]

    if intencion in ("archivos", "api", "voz", "general") and acciones:
        acciones = [a for a in acciones if a.get("accion") not in ("obtener_ip_local", "obtener_ip_publica")]

    resultados = ejecutar_acciones(acciones)
    hubo_error = any(_resultado_es_fallido(r) for r in resultados)

    if acciones and not hubo_error:
        guardar_habilidad(solicitud_original, acciones)
    registrar_leccion(solicitud_original, acciones, "\n".join(resultados), not hubo_error)

    respuesta_final = respuesta or "Sin respuesta de texto."
    if resultados:
        respuesta_final += "\n\n" + "\n".join(resultados)

    fallo_respuesta = any(p in respuesta_final.lower() for p in NEGATIVE_PATTERNS) or "no encontré un plan ejecutable" in respuesta_final.lower()
    if requiere_plan and (hubo_error or (not acciones and fallo_respuesta)):
        mejora_id, tipo, titulo_mejora, detalle_mejora, _ = crear_mejora_segura_desde_solicitud(
            solicitud_original,
            motivo="fallo tras reintentos automáticos",
            origen="fallo",
        )
        respuesta_final = construir_respuesta_respaldo(
            solicitud_original,
            motivo=detalle_mejora,
            mejora_id=mejora_id,
        )
        respuesta_final += (
            f"\n\nRegistré una mejora segura #{mejora_id} ({tipo}: {titulo_mejora}) para que puedas revisarla en el panel `Mejoras`."
        )

    guardar_mensaje_db("assistant", respuesta_final)
    registrar_ejecucion(solicitud_original, respuesta_final)
    return respuesta_final


class AgentCore:
    def __init__(self):
        init_db()
        asegurar_habilidades_base()
        depurar_datos_locales()
        limpiar_habilidades_y_automatizaciones_conflictivas()

    def process_prompt(self, prompt):
        return procesar_prompt_sync(prompt)


def _procesar_mensaje(texto):
    global ULTIMA_SOLICITUD_USUARIO
    chat_window.after(0, lambda: estado_var.set("Pensando..."))
    chat_window.after(0, expandir_si_compacto)
    texto_l = texto.lower()
    solicitud_original = texto.strip()
    tipo_tiempo = detectar_consulta_fecha_hora(solicitud_original)
    intencion = detectar_intencion_principal(solicitud_original)

    if not any(k in texto_l for k in ("aprende", "aprendelo", "apréndelo")):
        ULTIMA_SOLICITUD_USUARIO = solicitud_original

    if tipo_tiempo:
        cache_key = f"fecha_hora::{tipo_tiempo}::{normalizar_texto_cache(solicitud_original)}"
        ttl_seg = 10 if tipo_tiempo in ("hora", "fecha_hora") else 24 * 3600
        respuesta_cache = obtener_cache_respuesta(cache_key, max_age_seconds=ttl_seg)
        respuesta_final = respuesta_cache or formatear_fecha_hora_bonita(tipo_tiempo)
        if not respuesta_cache:
            guardar_cache_respuesta(cache_key, respuesta_final)
        def _post_tiempo():
            mensaje(respuesta_final, "ia")
            guardar_mensaje_db("assistant", respuesta_final)
            registrar_ejecucion(solicitud_original, respuesta_final)
            estado_var.set("Listo")
            if CONFIG.get("voz_activa", True):
                hablar_garantizado(respuesta_final)
        chat_window.after(0, _post_tiempo)
        return

    # Respuestas conversacionales locales para evitar bucles del planner.
    if intencion == "general":
        resp_local = respuesta_conversacional_local(solicitud_original)
        if resp_local:
            def _post_conv():
                mensaje(resp_local, "ia")
                guardar_mensaje_db("assistant", resp_local)
                registrar_ejecucion(solicitud_original, resp_local)
                estado_var.set("Listo")
                if CONFIG.get("voz_activa", True):
                    hablar_garantizado(resp_local)
            chat_window.after(0, _post_conv)
            return

    if any(k in texto_l for k in ("aprende a hacerlo", "aprende hacerlo", "aprendelo", "apréndelo")):
        if not ULTIMA_SOLICITUD_USUARIO:
            respuesta_final = "No tengo una solicitud previa para aprender. Dame primero una tarea concreta."
        else:
            acciones_aprendidas = acciones_locales_desde_texto(ULTIMA_SOLICITUD_USUARIO)
            if not acciones_aprendidas:
                respuesta_final = (
                    "Intenté aprender, pero no encontré un patrón automático para esa tarea todavía. "
                    "Dímela con más detalle (ruta, archivo o patrón)."
                )
            else:
                guardar_habilidad(ULTIMA_SOLICITUD_USUARIO, acciones_aprendidas)
                resultados_aprendidos = ejecutar_acciones(acciones_aprendidas)
                respuesta_final = (
                    "Aprendido y ejecutado en tiempo real.\n\n" + "\n".join(resultados_aprendidos)
                )

    # Ejecuta atajos locales solo en intenciones operativas explícitas.
    acciones_locales = acciones_locales_desde_texto(texto) if es_intencion_operativa(solicitud_original) else []
    if acciones_locales:
        acciones_locales = filtrar_acciones_por_intencion(acciones_locales, solicitud_original)
    if acciones_locales:
        resultados_locales = ejecutar_acciones(acciones_locales)
        exito_local = not any(
            "error" in (r or "").lower() or "no pudo" in (r or "").lower() or "no se pudo" in (r or "").lower()
            for r in resultados_locales
        )
        registrar_leccion(solicitud_original, acciones_locales, "\n".join(resultados_locales), exito_local)
        if exito_local:
            guardar_habilidad(solicitud_original, acciones_locales)
        respuesta_final = "\n".join(resultados_locales)

        def _post_local():
            mensaje(respuesta_final, "ia")
            guardar_mensaje_db("assistant", respuesta_final)
            registrar_ejecucion(solicitud_original, respuesta_final)
            estado_var.set("Listo")
            if CONFIG.get("voz_activa", True):
                hablar_garantizado(respuesta_final)

        chat_window.after(0, _post_local)
        return

    if es_intencion_operativa(solicitud_original):
        propuesta_auto = []
        t_auto = solicitud_original.lower()
        base_auto = resolver_ruta_alias(solicitud_original) or descubrir_ruta_escritorio()
        if ("crear" in t_auto or "genera" in t_auto) and ("docx" in t_auto or "word" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "documento_generado.docx"), "tipo": "docx", "contenido": "Documento generado automáticamente"}}]
        elif ("crear" in t_auto or "genera" in t_auto) and ("xlsx" in t_auto or "excel" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "tabla_generada.xlsx"), "tipo": "xlsx", "contenido": "columna1,columna2\nvalor1,valor2"}}]
        elif ("crear" in t_auto or "genera" in t_auto) and ("pptx" in t_auto or "powerpoint" in t_auto):
            propuesta_auto = [{"accion": "crear_archivo_especifico", "args": {"path": str(base_auto / "presentacion_generada.pptx"), "tipo": "pptx", "contenido": "Presentación generada automáticamente"}}]

        if propuesta_auto:
            resultados_auto = ejecutar_acciones(propuesta_auto)
            exito_auto = not any("error" in (r or "").lower() or "no pude" in (r or "").lower() for r in resultados_auto)
            registrar_leccion(solicitud_original, propuesta_auto, "\n".join(resultados_auto), exito_auto)
            if exito_auto:
                guardar_habilidad(solicitud_original, propuesta_auto)
            respuesta_auto = "\n".join(resultados_auto)

            def _post_auto():
                mensaje(respuesta_auto, "ia")
                guardar_mensaje_db("assistant", respuesta_auto)
                registrar_ejecucion(solicitud_original, respuesta_auto)
                estado_var.set("Listo")
                if CONFIG.get("voz_activa", True):
                    hablar_garantizado(respuesta_auto)
            chat_window.after(0, _post_auto)
            return

    try:
        respuesta_cruda = decidir(texto)
    except Exception as ex:
        respuesta_cruda = json.dumps(
            {
                "respuesta": f"Error consultando proveedor IA: {ex}",
                "acciones": [],
            },
            ensure_ascii=False,
        )

    try:
        data = json.loads(respuesta_cruda)
    except json.JSONDecodeError as ex:
        texto_plano = (respuesta_cruda or "").strip()
        if texto_plano:
            respuesta = texto_plano
        else:
            respuesta = "No devolviste JSON útil."
        data = {"respuesta": respuesta, "acciones": []}

    respuesta = data.get("respuesta", "").strip()
    acciones = data.get("acciones", [])
    acciones = filtrar_acciones_por_intencion(acciones, solicitud_original)
    if not acciones and any(
        s in respuesta.lower()
        for s in (
            "no puedo acceder",
            "no tengo la capacidad",
            "no puedo acceder directamente",
        )
    ):
        ruta_alias = resolver_ruta_alias(texto_l)
        if ruta_alias:
            acciones = [{"accion": "listar_directorio", "args": {"path": str(ruta_alias)}}]
            if not respuesta:
                respuesta = "Lo consulto localmente ahora mismo."

    # Solo activar planner cuando la solicitud parezca operativa.
    # Evita arrastrar acciones del turno previo en preguntas conversacionales.
    requiere_plan = es_intencion_operativa(solicitud_original)
    if requiere_plan and ((not acciones) or any(p in respuesta.lower() for p in NEGATIVE_PATTERNS)):
        try:
            r_plan, a_plan = generar_acciones_con_modelo(texto)
            if a_plan:
                acciones = a_plan
                if r_plan:
                    respuesta = r_plan
            else:
                # Si no hay plan válido, responde sin ejecutar comandos arbitrarios.
                acciones = []
                if not respuesta:
                    respuesta = (
                        "No encontré un plan ejecutable para esa petición aún. "
                        "Si quieres, te muestro opciones o me das una instrucción más concreta."
                    )
        except Exception as ex:
            acciones = []
            if not respuesta:
                respuesta = f"Error en planificación: {ex}"

    if not tipo_tiempo and acciones:
        acciones = [a for a in acciones if a.get("accion") != "obtener_hora"]

    if intencion in ("archivos", "api", "voz", "general") and acciones:
        acciones = [a for a in acciones if a.get("accion") not in ("obtener_ip_local", "obtener_ip_publica")]

    resultados = ejecutar_acciones(acciones)

    def _resultado_es_fallido(linea):
        s = (linea or "").lower()
        patrones = (
            "error",
            "no pude",
            "no tengo",
            "ruta no encontrada",
            "url inválida",
            "api.example.com",
            "failed",
            "max retries exceeded",
        )
        return any(p in s for p in patrones)

    hubo_error = any(_resultado_es_fallido(r) for r in resultados)

    # Aprende automáticamente lo que funcionó.
    if acciones and not hubo_error:
        guardar_habilidad(solicitud_original, acciones)
    registrar_leccion(solicitud_original, acciones, "\n".join(resultados), not hubo_error)

    respuesta_final = respuesta or "Sin respuesta de texto."
    if resultados:
        respuesta_final += "\n\n" + "\n".join(resultados)

    def _post():
        mensaje(respuesta_final, "ia")
        guardar_mensaje_db("assistant", respuesta_final)
        registrar_ejecucion(solicitud_original, respuesta_final)
        estado_var.set("Listo")
        if CONFIG.get("voz_activa", True):
            hablar_garantizado(respuesta_final)

    chat_window.after(0, _post)


def enviar(event=None):
    expandir_si_compacto()
    texto = obtener_texto_entrada()
    limpiar_entrada_ui()
    if not texto:
        return "break" if event is not None else None

    mensaje(texto, "user")
    guardar_mensaje_db("user", texto)
    estado_var.set("Pensando...")
    def _worker():
        try:
            _procesar_mensaje(texto)
        except Exception as ex:
            err = f"Error interno procesando mensaje: {ex}"
            def _post_err():
                mensaje(err, "ia")
                guardar_mensaje_db("assistant", err)
                registrar_ejecucion(texto, err)
                estado_var.set("Listo")
            chat_window.after(0, _post_err)
    threading.Thread(target=_worker, daemon=True).start()
    return "break" if event is not None else None


def escuchar_y_enviar():
    global escucha_activa
    if sr is None:
        estado_var.set("Falta SpeechRecognition en este entorno.")
        chat_window.after(2200, lambda: estado_var.set("Listo"))
        return
    if escucha_activa:
        return
    if not CONFIG.get("voz_entrada_activa", True):
        estado_var.set("Micrófono desactivado")
        chat_window.after(1500, lambda: estado_var.set("Listo"))
        return

    escucha_activa = True
    estado_var.set("Escuchando...")
    btn_micro.configure(text="Escuchando...", state="disabled")

    def _run():
        global escucha_activa
        r = sr.Recognizer()
        texto = ""
        error = None
        try:
            mic_nombre = mic_var.get().strip()
            mic_index = None
            nombres = listar_microfonos()
            if mic_nombre and mic_nombre in nombres:
                mic_index = nombres.index(mic_nombre)
            if mic_nombre and mic_index is None:
                raise RuntimeError(f"Micrófono no encontrado: {mic_nombre}")

            with sr.Microphone(device_index=mic_index) as source:
                r.adjust_for_ambient_noise(source, duration=0.5)
                audio = r.listen(source, timeout=8, phrase_time_limit=14)
            idioma = CONFIG.get("voz_entrada_idioma", "es-ES")
            texto = r.recognize_google(audio, language=idioma)
        except sr.WaitTimeoutError:
            error = "No detecté voz a tiempo."
        except sr.UnknownValueError:
            error = "No pude entender lo que dijiste."
        except sr.RequestError as ex:
            error = f"Error de reconocimiento: {ex}"
        except Exception as ex:
            error = f"Error con micrófono: {ex}"

        def _post():
            global escucha_activa
            escucha_activa = False
            btn_micro.configure(text="Hablar", state="normal")
            if error:
                estado_var.set(error)
                chat_window.after(2200, lambda: estado_var.set("Listo"))
                return

            poner_texto_entrada(texto)
            estado_var.set(f"Texto reconocido ({mic_var.get() or 'predeterminado'})")
            chat_window.after(1200, lambda: estado_var.set("Listo"))
            enviar()

        chat_window.after(0, _post)

    threading.Thread(target=_run, daemon=True).start()

def refrescar_microfonos():
    nombres = listar_microfonos()
    combo_mic.configure(values=nombres, state="readonly" if nombres else "disabled")
    if nombres and mic_var.get() not in nombres:
        mic_var.set(nombres[0])
    CONFIG["voz_entrada_microfono"] = mic_var.get()
    guardar_config(CONFIG)
    estado_var.set("Micrófonos actualizados")
    chat_window.after(1200, lambda: estado_var.set("Listo"))
def on_mic_change(_event=None):
    CONFIG["voz_entrada_microfono"] = mic_var.get()
    guardar_config(CONFIG)


def configurar_interfaz():
    global root
    global chat_window
    global style
    global estado_var
    global header
    global titulo
    global hero_subtitle_var
    global provider_badge_var
    global provider_badge
    global tts_status_var
    global tts_estado
    global model_var
    global combo_model
    global provider_var
    global combo_provider
    global online_model_var
    global combo_online_model
    global mic_var
    global combo_mic
    global voz_style_var
    global combo_voz_style
    global voz_speed_var
    global combo_voz_speed
    global voz_var
    global chk_voz
    global skills_auto_var
    global chk_skills_auto
    global btn_compacto
    global btn_settings
    global btn_guardar
    global btn_habilidades
    global btn_mejoras
    global btn_automatizaciones
    global btn_integraciones
    global btn_key
    global frame_chat
    global canvas
    global scrollable_frame
    global settings_panel
    global settings_visible
    global frame_input
    global entrada
    global btn_limpiar
    global btn_micro
    global btn_refresh_mic
    global btn_enviar
    global chat_compacto
    global auto_compact_job

    if chat_window is not None:
        return chat_window

    root = tk.Tk()
    chat_window = root
    chat_window.title("Agente IA híbrido")

    size_actual = normalizar_window_size(CONFIG.get("window_size"))
    m_size = re.match(r"^(\d+)x(\d+)$", size_actual)
    if m_size:
        ancho = max(760, int(m_size.group(1)))
        alto = max(720, int(m_size.group(2)))
        size_actual = f"{ancho}x{alto}"

    chat_window.geometry(
        f"{size_actual}+{CONFIG.get('window_x', 1000)}+{CONFIG.get('window_y', 140)}"
    )
    chat_window.configure(bg="#08111d")
    chat_window.minsize(560, 620)

    style = ttk.Style()
    style.theme_use("clam")
    style.configure(
        "TCombobox",
        fieldbackground="#08111d",
        foreground="#e2e8f0",
        arrowsize=13,
        borderwidth=0,
        lightcolor="#111827",
        darkcolor="#111827",
    )
    style.configure(
        "Vertical.TScrollbar",
        background="#1e293b",
        troughcolor="#0b1220",
        bordercolor="#0b1220",
        arrowcolor="#cbd5e1",
    )
    style.configure(
        "Treeview",
        background="#0b1220",
        foreground="#e2e8f0",
        fieldbackground="#0b1220",
        rowheight=26,
    )
    style.configure(
        "Treeview.Heading",
        background="#111827",
        foreground="#e2e8f0",
        relief="flat",
    )

    chat_compacto = False
    auto_compact_job = None
    settings_visible = True

    def make_button(parent, text, command, variant="ghost", padx=12):
        palette = {
            "ghost": {"bg": "#111827", "fg": "#cbd5e1"},
            "primary": {"bg": "#2563eb", "fg": "#eff6ff"},
            "accent": {"bg": "#0f766e", "fg": "#ecfeff"},
        }
        colors = palette.get(variant, palette["ghost"])
        return tk.Button(
            parent,
            text=text,
            command=command,
            bg=colors["bg"],
            fg=colors["fg"],
            activebackground=colors["bg"],
            activeforeground=colors["fg"],
            relief="flat",
            bd=0,
            padx=padx,
            pady=7,
            font=("Segoe UI Semibold", 9),
            cursor="hand2",
        )

    shell = tk.Frame(chat_window, bg="#08111d")
    shell.pack(fill="both", expand=True, padx=16, pady=16)

    header = tk.Frame(shell, bg="#0f172a", padx=18, pady=18)
    header.pack(fill="x", pady=(0, 14))

    hero_top = tk.Frame(header, bg="#0f172a")
    hero_top.pack(fill="x")

    left_col = tk.Frame(hero_top, bg="#0f172a")
    left_col.pack(side="left", fill="x", expand=True)

    titulo = tk.Label(
        left_col,
        text="Agente IA híbrido",
        bg="#0f172a",
        fg="#f8fafc",
        font=("Segoe UI Semibold", 16),
    )
    titulo.pack(anchor="w")

    hero_subtitle_var = tk.StringVar(value="Modo manual activado.")
    tk.Label(
        left_col,
        textvariable=hero_subtitle_var,
        bg="#0f172a",
        fg="#94a3b8",
        justify="left",
        wraplength=520,
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(4, 0))

    right_col = tk.Frame(hero_top, bg="#0f172a")
    right_col.pack(side="right", anchor="ne")

    provider_badge_var = tk.StringVar(value="LOCAL")
    provider_badge = tk.Label(
        right_col,
        textvariable=provider_badge_var,
        bg="#163323",
        fg="#86efac",
        padx=10,
        pady=5,
        font=("Segoe UI Semibold", 9),
    )
    provider_badge.pack(anchor="e")

    estado_var = tk.StringVar(value="Listo")
    tk.Label(
        right_col,
        textvariable=estado_var,
        bg="#111827",
        fg="#e2e8f0",
        padx=10,
        pady=5,
        font=("Segoe UI Semibold", 9),
    ).pack(anchor="e", pady=(8, 4))

    tts_status_var = tk.StringVar(value="🛑 silencio")
    tts_estado = tk.Label(
        right_col,
        textvariable=tts_status_var,
        bg="#0f172a",
        fg="#60a5fa",
        font=("Segoe UI", 9),
    )
    tts_estado.pack(anchor="e")

    action_row = tk.Frame(header, bg="#0f172a")
    action_row.pack(fill="x", pady=(16, 0))

    btn_settings = make_button(action_row, "Ocultar ajustes", alternar_panel_ajustes)
    btn_settings.pack(side="left")

    btn_guardar = make_button(action_row, "Guardar", guardar_preferencias, variant="primary")
    btn_guardar.pack(side="right")

    btn_compacto = make_button(action_row, "Compacto", alternar_compacto, padx=10)
    btn_compacto.pack(side="right", padx=(0, 8))

    settings_panel = tk.Frame(shell, bg="#111827", padx=18, pady=18)
    settings_panel.pack(fill="x", pady=(0, 14))

    tk.Label(
        settings_panel,
        text="Control manual de ejecución",
        bg="#111827",
        fg="#f8fafc",
        font=("Segoe UI Semibold", 11),
    ).pack(anchor="w")
    tk.Label(
        settings_panel,
        text="Aquí eliges tú si el agente trabaja con IA local o con API. El proveedor activo no cambia solo.",
        bg="#111827",
        fg="#94a3b8",
        wraplength=760,
        justify="left",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(4, 12))

    row_provider = tk.Frame(settings_panel, bg="#111827")
    row_provider.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_provider,
        text="Proveedor",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    provider_var = tk.StringVar(value=proveedor_ui_value(CONFIG.get("proveedor_ia", "local")))
    combo_provider = ttk.Combobox(
        row_provider,
        textvariable=provider_var,
        values=["local", "api"],
        width=12,
        state="readonly",
    )
    combo_provider.pack(side="left")

    row_model_local = tk.Frame(settings_panel, bg="#111827")
    row_model_local.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_model_local,
        text="Modelo local",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    model_var = tk.StringVar(value=CONFIG["model"])
    combo_model = ttk.Combobox(
        row_model_local,
        textvariable=model_var,
        values=["qwen2.5:7b", "llama3.1:8b", "mistral:7b", "phi3:mini"],
        width=24,
        state="readonly",
    )
    combo_model.pack(side="left", fill="x", expand=True)

    row_model_api = tk.Frame(settings_panel, bg="#111827")
    row_model_api.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_model_api,
        text="Modelo API",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    online_model_var = tk.StringVar(value=CONFIG.get("modelo_online", "llama-3.1-8b-instant"))
    combo_online_model = ttk.Combobox(
        row_model_api,
        textvariable=online_model_var,
        values=["llama-3.1-8b-instant", "llama-3.3-70b-versatile", "mixtral-8x7b-32768"],
        width=30,
        state="readonly",
    )
    combo_online_model.pack(side="left", fill="x", expand=True)

    row_voice = tk.Frame(settings_panel, bg="#111827")
    row_voice.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_voice,
        text="Voz",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    voz_var = tk.BooleanVar(value=CONFIG["voz_activa"])
    chk_voz = tk.Checkbutton(
        row_voice,
        text="Activada",
        variable=voz_var,
        command=on_toggle_voz,
        bg="#111827",
        fg="#cbd5e1",
        selectcolor="#1f2937",
        activebackground="#111827",
        activeforeground="#cbd5e1",
        font=("Segoe UI", 9),
    )
    chk_voz.pack(side="left", padx=(0, 10))

    voz_style_var = tk.StringVar(value=CONFIG.get("voz_style", "Natural"))
    combo_voz_style = ttk.Combobox(
        row_voice,
        textvariable=voz_style_var,
        values=["Natural", "Suave", "Profunda", "Femenina", "Masculina"],
        width=14,
        state="readonly",
    )
    combo_voz_style.pack(side="left", padx=(0, 8))

    voz_speed_var = tk.StringVar(value=CONFIG.get("voz_speed_label", "Normal"))
    combo_voz_speed = ttk.Combobox(
        row_voice,
        textvariable=voz_speed_var,
        values=["Lenta", "Normal", "Rápida"],
        width=12,
        state="readonly",
    )
    combo_voz_speed.pack(side="left")

    row_mic = tk.Frame(settings_panel, bg="#111827")
    row_mic.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_mic,
        text="Micrófono",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    mic_names = listar_microfonos()
    mic_var = tk.StringVar(value=CONFIG.get("voz_entrada_microfono", ""))
    combo_mic = ttk.Combobox(
        row_mic,
        textvariable=mic_var,
        values=mic_names,
        width=34,
        state="readonly" if mic_names else "disabled",
    )
    if not mic_var.get() and mic_names:
        mic_var.set(mic_names[0])
    combo_mic.pack(side="left", fill="x", expand=True)

    btn_refresh_mic = make_button(row_mic, "Actualizar", refrescar_microfonos, padx=10)
    btn_refresh_mic.pack(side="left", padx=(8, 0))

    row_learning = tk.Frame(settings_panel, bg="#111827")
    row_learning.pack(fill="x", pady=(0, 10))
    tk.Label(
        row_learning,
        text="Atajos",
        width=16,
        anchor="w",
        bg="#111827",
        fg="#cbd5e1",
        font=("Segoe UI Semibold", 9),
    ).pack(side="left")
    skills_auto_var = tk.BooleanVar(value=bool(CONFIG.get("usar_habilidades_auto", False)))
    chk_skills_auto = tk.Checkbutton(
        row_learning,
        text="Aprender ejecuciones exitosas",
        variable=skills_auto_var,
        command=on_toggle_skills_auto,
        bg="#111827",
        fg="#cbd5e1",
        selectcolor="#1f2937",
        activebackground="#111827",
        activeforeground="#cbd5e1",
        font=("Segoe UI", 9),
    )
    chk_skills_auto.pack(side="left")

    tk.Label(
        settings_panel,
        text="Gestión rápida",
        bg="#111827",
        fg="#94a3b8",
        font=("Segoe UI Semibold", 9),
    ).pack(anchor="w", pady=(8, 6))

    tools_row = tk.Frame(settings_panel, bg="#111827")
    tools_row.pack(fill="x")
    btn_habilidades = make_button(tools_row, "Habilidades", abrir_panel_habilidades)
    btn_habilidades.pack(side="left")
    btn_mejoras = make_button(tools_row, "Mejoras", abrir_panel_mejoras)
    btn_mejoras.pack(side="left", padx=(8, 0))
    btn_automatizaciones = make_button(tools_row, "Automatizaciones", abrir_panel_automatizaciones)
    btn_automatizaciones.pack(side="left", padx=(8, 0))
    btn_integraciones = make_button(tools_row, "Integraciones", abrir_panel_integraciones)
    btn_integraciones.pack(side="left", padx=(8, 0))
    btn_key = make_button(tools_row, "Groq Key", pedir_api_key_groq, variant="accent")
    btn_key.pack(side="left", padx=(8, 0))

    chat_card = tk.Frame(shell, bg="#0b1220", padx=12, pady=12)
    chat_card.pack(fill="both", expand=True, pady=(0, 14))

    chat_head = tk.Frame(chat_card, bg="#0b1220")
    chat_head.pack(fill="x", padx=6, pady=(4, 8))
    tk.Label(
        chat_head,
        text="Conversación",
        bg="#0b1220",
        fg="#f8fafc",
        font=("Segoe UI Semibold", 11),
    ).pack(anchor="w")
    tk.Label(
        chat_head,
        text="El proveedor se mantiene fijo hasta que cambies el selector manual.",
        bg="#0b1220",
        fg="#64748b",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(3, 0))

    frame_chat = tk.Frame(chat_card, bg="#08111d")
    frame_chat.pack(fill="both", expand=True)

    canvas = tk.Canvas(frame_chat, bg="#08111d", highlightthickness=0, bd=0)
    scrollbar = ttk.Scrollbar(
        frame_chat, orient="vertical", style="Vertical.TScrollbar", command=canvas.yview
    )
    scrollable_frame = tk.Frame(canvas, bg="#08111d")
    scroll_window = canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")

    def _sync_scrollregion(_event=None):
        canvas.configure(scrollregion=canvas.bbox("all"))

    def _sync_width(_event=None):
        canvas.itemconfigure(scroll_window, width=canvas.winfo_width())

    scrollable_frame.bind("<Configure>", _sync_scrollregion)
    canvas.bind("<Configure>", _sync_width)
    canvas.configure(yscrollcommand=scrollbar.set)
    canvas.pack(side="left", fill="both", expand=True)
    scrollbar.pack(side="right", fill="y")

    frame_input = tk.Frame(shell, bg="#0f172a", padx=16, pady=16)
    frame_input.pack(fill="x")

    tk.Label(
        frame_input,
        text="Escribe una instrucción para tu agente",
        bg="#0f172a",
        fg="#f8fafc",
        font=("Segoe UI Semibold", 11),
    ).pack(anchor="w")
    tk.Label(
        frame_input,
        text="Enter envía. Shift+Enter inserta una nueva línea. El modo local/api siempre lo eliges tú.",
        bg="#0f172a",
        fg="#94a3b8",
        font=("Segoe UI", 9),
    ).pack(anchor="w", pady=(3, 8))

    entrada = tk.Text(
        frame_input,
        height=3,
        wrap="word",
        bg="#08111d",
        fg="#e2e8f0",
        insertbackground="#e2e8f0",
        relief="flat",
        bd=0,
        padx=12,
        pady=10,
        font=("Segoe UI", 11),
    )
    entrada.pack(fill="x")

    composer_actions = tk.Frame(frame_input, bg="#0f172a")
    composer_actions.pack(fill="x", pady=(10, 0))

    btn_limpiar = make_button(composer_actions, "Limpiar", limpiar_historial, padx=10)
    btn_limpiar.pack(side="left")

    btn_micro = make_button(composer_actions, "Hablar", escuchar_y_enviar, padx=10)
    btn_micro.pack(side="left", padx=(8, 0))

    btn_enviar = make_button(composer_actions, "Enviar", enviar, variant="primary", padx=16)
    btn_enviar.pack(side="right")

    entrada.bind("<Return>", on_textbox_return)
    entrada.bind("<Shift-Return>", on_textbox_shift_return)
    combo_model.bind("<<ComboboxSelected>>", on_model_change)
    combo_voz_style.bind("<<ComboboxSelected>>", on_voice_ui_change)
    combo_voz_speed.bind("<<ComboboxSelected>>", on_voice_ui_change)
    combo_provider.bind("<<ComboboxSelected>>", on_provider_change)
    combo_online_model.bind("<<ComboboxSelected>>", on_online_model_change)
    combo_mic.bind("<<ComboboxSelected>>", on_mic_change)
    chat_window.bind("<FocusIn>", lambda e: expandir_si_compacto())
    chat_window.bind("<Configure>", lambda e: guardar_geometria_actual())
    chat_window.protocol("WM_DELETE_WINDOW", cerrar_app)

    actualizar_resumen_visual()
    entrada.focus_set()

    return chat_window


def depurar_datos_locales():
    return


def limpiar_habilidades_y_automatizaciones_conflictivas():
    with sqlite3.connect(DB_PATH) as conn:
        rows = conn.execute("SELECT id, trigger_texto, acciones_json FROM habilidades ORDER BY id ASC").fetchall()
        to_delete = []
        for hid, trigger_texto, acciones_json in rows:
            t = normalizar_texto_cache(trigger_texto or "")
            if not t:
                continue
            try:
                acciones = json.loads(acciones_json or "[]")
            except Exception:
                to_delete.append(hid)
                continue
            if not isinstance(acciones, list) or not acciones:
                to_delete.append(hid)
                continue

            # Reglas de saneamiento para evitar recuerdos tóxicos.
            tiene_archivos = ("archivo" in t or "archivos" in t)
            tiene_carpeta = ("carpeta" in t)
            primer_accion = str(acciones[0].get("accion", "")).strip() if isinstance(acciones[0], dict) else ""

            # Si el trigger pide archivos, no debe aprender acciones de abrir_ruta/crear_carpeta.
            if tiene_archivos and primer_accion in ("abrir_ruta", "crear_carpeta"):
                to_delete.append(hid)
                continue
            # Si dice "carpeta llamada X", no debe guardar literal "llamada x" en el path.
            if tiene_carpeta and "llamada " in t:
                path_bad = False
                for a in acciones:
                    if isinstance(a, dict):
                        p = str((a.get("args") or {}).get("path", "")).lower()
                        if "\\llamada " in p or "/llamada " in p:
                            path_bad = True
                            break
                if path_bad:
                    to_delete.append(hid)
                    continue

        if to_delete:
            q = ",".join("?" for _ in to_delete)
            conn.execute(f"DELETE FROM habilidades WHERE id IN ({q})", tuple(to_delete))
            conn.commit()


init_db()
asegurar_habilidades_base()
depurar_datos_locales()
limpiar_habilidades_y_automatizaciones_conflictivas()

if __name__ == "__main__":
    configurar_interfaz()
    mensaje("Asistente listo. Ventana minimalista activa.", "system")
    programar_auto_compacto()
    root.mainloop()


